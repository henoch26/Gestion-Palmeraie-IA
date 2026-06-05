"""
Script de génération de la présentation PowerPoint
Projet : Système intelligent de gestion du palmier à huile
Auteur : ABO KOKOUTSE ADZEYI HENOCH
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ─────────────────── Palette de couleurs ───────────────────
BLEU_FONCE   = RGBColor(0x1A, 0x3A, 0x5C)   # Titres principaux
VERT_FONCE   = RGBColor(0x1B, 0x5E, 0x20)   # Accents agriculture
VERT_CLAIR   = RGBColor(0x43, 0xA0, 0x47)   # Sous-titres / badges
OR           = RGBColor(0xF9, 0xA8, 0x25)   # Highlights / palmier
BLANC        = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR   = RGBColor(0xF5, 0xF5, 0xF5)
GRIS_TEXTE   = RGBColor(0x42, 0x42, 0x42)
BLEU_CIEL    = RGBColor(0xE3, 0xF2, 0xFD)   # Fond header tableau
VERT_PALE    = RGBColor(0xE8, 0xF5, 0xE9)   # Fond lignes paires

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]   # Entièrement vide

# ─────────────────── Helpers ────────────────────────────────

def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=12, bold=False, italic=False,
                 color=GRIS_TEXTE, align=PP_ALIGN.LEFT,
                 v_anchor=None, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_paragraph(tf, text, font_size=11, bold=False, italic=False,
                  color=GRIS_TEXTE, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None):
    """Bandeau supérieur bleu foncé avec titre et barre verte"""
    # Fond bandeau
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), BLEU_FONCE)
    # Barre accent verte en bas du bandeau
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.065), VERT_CLAIR)
    # Titre
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.65),
                 font_size=26, bold=True, color=BLANC, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.38),
                     font_size=13, bold=False, italic=True, color=OR, align=PP_ALIGN.LEFT)

def slide_footer(slide, page_num, total=20):
    """Pied de page sobre"""
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), BLEU_FONCE)
    add_text_box(slide, "UFHB · ABO KOKOUTSE ADZEYI HENOCH · 2025-2026",
                 Inches(0.3), Inches(7.16), Inches(10), Inches(0.3),
                 font_size=8, color=BLANC, align=PP_ALIGN.LEFT)
    add_text_box(slide, f"{page_num} / {total}",
                 Inches(12.3), Inches(7.16), Inches(0.8), Inches(0.3),
                 font_size=8, bold=True, color=OR, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, x, y, w, h, title=None,
               title_color=BLEU_FONCE, bullet="▸", font_size=11):
    """Boîte avec liste à puces"""
    if title:
        add_text_box(slide, title, x, y, w, Inches(0.35),
                     font_size=13, bold=True, color=title_color)
        y += Inches(0.38)
        h -= Inches(0.38)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = GRIS_TEXTE

def add_table(slide, headers, rows, x, y, w, h,
              header_fill=BLEU_FONCE, row_alt=VERT_PALE,
              font_size=9, header_font_size=9):
    """Tableau stylisé"""
    col_count = len(headers)
    row_count = len(rows) + 1
    table = slide.shapes.add_table(row_count, col_count, x, y, w, h).table

    col_w = w // col_count
    for i in range(col_count):
        table.columns[i].width = col_w

    # En-têtes
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(header_font_size)
        run.font.color.rgb = BLANC

    # Données
    for ri, row_data in enumerate(rows):
        fill_color = row_alt if ri % 2 == 0 else BLANC
        for ci, val in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.color.rgb = GRIS_TEXTE

    return table

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 1 – PAGE DE TITRE
# ═══════════════════════════════════════════════════════════════════
def slide_01_titre():
    slide = prs.slides.add_slide(BLANK)

    # Fond dégradé simulé : 2 rectangles
    add_rect(slide, 0, 0, SLIDE_W, Inches(4.2), BLEU_FONCE)
    add_rect(slide, 0, Inches(4.2), SLIDE_W, Inches(3.3), VERT_PALE)

    # Barre dorée horizontale séparatrice
    add_rect(slide, 0, Inches(4.1), SLIDE_W, Inches(0.12), OR)

    # Ligne décorative gauche
    add_rect(slide, Inches(0.35), Inches(0.5), Inches(0.08), Inches(3.3), OR)

    # Titre principal
    add_text_box(slide,
                 "Système Intelligent de Gestion Automatique\nde la Récolte et de la Productivité\ndu Palmier à Huile",
                 Inches(0.65), Inches(0.55), Inches(11.8), Inches(2.5),
                 font_size=28, bold=True, color=BLANC, align=PP_ALIGN.LEFT)

    # Sous-titre
    add_text_box(slide, "Application Web de Gestion de Palmeraie",
                 Inches(0.65), Inches(3.1), Inches(11.5), Inches(0.6),
                 font_size=16, italic=True, color=OR, align=PP_ALIGN.LEFT)

    # Infos étudiant
    add_text_box(slide, "Étudiant  :  ABO KOKOUTSE ADZEYI HENOCH",
                 Inches(0.65), Inches(4.55), Inches(8), Inches(0.45),
                 font_size=14, bold=True, color=BLEU_FONCE, align=PP_ALIGN.LEFT)

    add_text_box(slide, "Université  :  UFHB – Université Félix Houphouët-Boigny · Abidjan",
                 Inches(0.65), Inches(5.05), Inches(10), Inches(0.4),
                 font_size=12, color=GRIS_TEXTE, align=PP_ALIGN.LEFT)

    add_text_box(slide, "Master 2 – Base de Données et Génie Logiciel (BDGL)",
                 Inches(0.65), Inches(5.5), Inches(10), Inches(0.4),
                 font_size=12, color=GRIS_TEXTE, align=PP_ALIGN.LEFT)

    add_text_box(slide, "Année académique 2025 – 2026",
                 Inches(0.65), Inches(5.95), Inches(8), Inches(0.4),
                 font_size=12, italic=True, color=VERT_CLAIR, align=PP_ALIGN.LEFT)

    # Logo palmier (emoji via texte)
    add_text_box(slide, "🌴", Inches(11.2), Inches(0.4), Inches(1.8), Inches(3.2),
                 font_size=80, align=PP_ALIGN.CENTER, color=OR)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 2 – PLAN
# ═══════════════════════════════════════════════════════════════════
def slide_02_plan():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Plan de la Présentation")
    slide_footer(slide, 2)

    sections = [
        ("I.",   "Contexte & Problématique"),
        ("II.",  "Objectifs du projet"),
        ("III.", "Architecture globale du système"),
        ("IV.",  "Processus de développement"),
        ("V.",   "Outils et technologies utilisés"),
        ("VI.",  "Modules de l'application"),
        ("VII.", "Schéma de la base de données"),
        ("VIII.","Détail des tables (attributs)"),
        ("IX.",  "Fonctionnalités principales"),
        ("X.",   "Tableau de bord & Analyses"),
        ("XI.",  "Conclusion & Perspectives"),
    ]

    col1 = sections[:6]
    col2 = sections[6:]

    for idx, (num, label) in enumerate(col1):
        y = Inches(1.45) + idx * Inches(0.82)
        add_rect(slide, Inches(0.4), y, Inches(0.55), Inches(0.52), BLEU_FONCE)
        add_text_box(slide, num, Inches(0.4), y, Inches(0.55), Inches(0.52),
                     font_size=11, bold=True, color=OR, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, Inches(1.05), y + Inches(0.06), Inches(5.8), Inches(0.42),
                     font_size=13, color=GRIS_TEXTE)

    for idx, (num, label) in enumerate(col2):
        y = Inches(1.45) + idx * Inches(0.82)
        add_rect(slide, Inches(7.0), y, Inches(0.55), Inches(0.52), VERT_FONCE)
        add_text_box(slide, num, Inches(7.0), y, Inches(0.55), Inches(0.52),
                     font_size=11, bold=True, color=OR, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, Inches(7.65), y + Inches(0.06), Inches(5.3), Inches(0.42),
                     font_size=13, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 3 – CONTEXTE & PROBLÉMATIQUE
# ═══════════════════════════════════════════════════════════════════
def slide_03_contexte():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Contexte & Problématique",
                 "Pourquoi ce projet ?")
    slide_footer(slide, 3)

    # Contexte
    add_text_box(slide, "Contexte", Inches(0.4), Inches(1.45), Inches(12), Inches(0.4),
                 font_size=15, bold=True, color=BLEU_FONCE)
    add_rect(slide, Inches(0.4), Inches(1.87), Inches(12.5), Inches(0.04), VERT_CLAIR)

    contexte_items = [
        "La filière palmier à huile est un secteur économique stratégique en Côte d'Ivoire",
        "Les exploitations produisent des volumes importants de régimes qui nécessitent un suivi rigoureux",
        "La gestion manuelle (papier, tableurs) génère des erreurs, des pertes de données et un manque de visibilité",
        "Le besoin d'un outil numérique centralisé s'impose pour moderniser la gestion de la palmeraie",
    ]
    bullet_box(slide, contexte_items, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.8),
               font_size=12, bullet="●")

    # Problématique
    add_text_box(slide, "Problématique", Inches(0.4), Inches(3.9), Inches(12), Inches(0.4),
                 font_size=15, bold=True, color=VERT_FONCE)
    add_rect(slide, Inches(0.4), Inches(4.32), Inches(12.5), Inches(0.04), OR)

    prob_items = [
        "Comment automatiser et centraliser la saisie, le suivi et l'analyse des données de récolte ?",
        "Comment permettre une traçabilité complète des activités (récolte, travaux, matériels) ?",
        "Comment fournir des tableaux de bord dynamiques pour aider à la prise de décision ?",
    ]
    bullet_box(slide, prob_items, Inches(0.5), Inches(4.45), Inches(12.3), Inches(1.5),
               font_size=12, bullet="?")

    # Encadré solution
    add_rect(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.7), BLEU_FONCE)
    add_text_box(slide,
                 "➜  Solution : Une application web full-stack combinant Django REST Framework (backend) et React (frontend) avec une base de données PostgreSQL",
                 Inches(0.6), Inches(6.18), Inches(12.1), Inches(0.55),
                 font_size=12, bold=True, color=BLANC)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 4 – OBJECTIFS
# ═══════════════════════════════════════════════════════════════════
def slide_04_objectifs():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Objectifs du Projet",
                 "Ce que l'application doit permettre")
    slide_footer(slide, 4)

    objectifs = [
        ("🌿", "Gestion des Récoltes",
         "Saisir et suivre les fiches de récolte quotidiennes : superviseurs, récolteurs, régimes (grands / moyens / petits), secteurs, ventes"),
        ("👷", "Gestion des Travaux",
         "Enregistrer les fiches de travaux agricoles avec la répartition des tâches, les consommables utilisés et les secteurs couverts"),
        ("🔧", "Inventaire des Matériels",
         "Suivre l'état physique et le statut d'utilisation de chaque équipement de la palmeraie"),
        ("👥", "Annuaire des Récolteurs",
         "Gérer un répertoire centralisé des récolteurs avec attribution automatique de codes uniques (REC-001, ...)"),
        ("🗺️", "Gestion des Secteurs",
         "Cartographier les secteurs (superficie, relief, type de sol) avec génération automatique de codes"),
        ("📊", "Tableau de Bord Analytique",
         "Visualiser les statistiques de production filtrées par année, secteur, récolteur et type de régime"),
    ]

    cols = [objectifs[:3], objectifs[3:]]
    for col_idx, col in enumerate(cols):
        x_base = Inches(0.3) + col_idx * Inches(6.5)
        for row_idx, (icon, titre, desc) in enumerate(col):
            y = Inches(1.45) + row_idx * Inches(1.85)
            add_rect(slide, x_base, y, Inches(6.2), Inches(1.7), BLANC,
                     line_color=VERT_CLAIR, line_width_pt=1)
            add_text_box(slide, icon, x_base + Inches(0.1), y + Inches(0.1),
                         Inches(0.7), Inches(0.65), font_size=22, align=PP_ALIGN.CENTER)
            add_text_box(slide, titre, x_base + Inches(0.85), y + Inches(0.1),
                         Inches(5.1), Inches(0.42),
                         font_size=13, bold=True, color=BLEU_FONCE)
            add_text_box(slide, desc, x_base + Inches(0.85), y + Inches(0.55),
                         Inches(5.1), Inches(1.1),
                         font_size=10, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 5 – ARCHITECTURE GLOBALE
# ═══════════════════════════════════════════════════════════════════
def slide_05_architecture():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Architecture Globale du Système",
                 "Architecture 3-tiers : Client · Serveur · Base de données")
    slide_footer(slide, 5)

    # Couche 1 – Frontend
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(3.5), Inches(4.8), BLEU_CIEL,
             line_color=BLEU_FONCE, line_width_pt=1.5)
    add_text_box(slide, "FRONTEND", Inches(0.4), Inches(1.5), Inches(3.5), Inches(0.5),
                 font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(3.5), Inches(0.45), BLEU_FONCE)
    items_front = ["React 19", "React Router DOM", "Axios (HTTP)", "Vite 7 (build)", "ESLint (qualité)"]
    bullet_box(slide, items_front, Inches(0.55), Inches(2.1), Inches(3.2), Inches(3.5),
               font_size=11, bullet="▸")
    add_text_box(slide, "Navigateur Web\n(Interface utilisateur)", Inches(0.5), Inches(5.8),
                 Inches(3.2), Inches(0.6), font_size=10, italic=True, color=BLEU_FONCE, align=PP_ALIGN.CENTER)

    # Flèche 1
    add_text_box(slide, "HTTP / REST API\n←→", Inches(3.95), Inches(3.5), Inches(1.5), Inches(0.8),
                 font_size=11, bold=True, color=OR, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.95), Inches(3.9), Inches(1.45), Inches(0.06), OR)

    # Couche 2 – Backend
    add_rect(slide, Inches(5.5), Inches(1.5), Inches(3.5), Inches(4.8), VERT_PALE,
             line_color=VERT_FONCE, line_width_pt=1.5)
    add_rect(slide, Inches(5.5), Inches(1.5), Inches(3.5), Inches(0.45), VERT_FONCE)
    add_text_box(slide, "BACKEND", Inches(5.5), Inches(1.5), Inches(3.5), Inches(0.5),
                 font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    items_back = ["Django 5.2", "Django REST Framework", "Token Authentication", "CORS Headers", "Modules : récoltes, travaux,\nmatériels, récolteurs, secteurs"]
    bullet_box(slide, items_back, Inches(5.65), Inches(2.1), Inches(3.2), Inches(3.5),
               font_size=11, bullet="▸")
    add_text_box(slide, "API RESTful\n(Logique métier)", Inches(5.6), Inches(5.8),
                 Inches(3.2), Inches(0.6), font_size=10, italic=True, color=VERT_FONCE, align=PP_ALIGN.CENTER)

    # Flèche 2
    add_text_box(slide, "ORM Django\n←→", Inches(9.05), Inches(3.5), Inches(1.5), Inches(0.8),
                 font_size=11, bold=True, color=OR, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(9.05), Inches(3.9), Inches(1.45), Inches(0.06), OR)

    # Couche 3 – BD
    add_rect(slide, Inches(10.6), Inches(1.5), Inches(2.4), Inches(4.8), RGBColor(0xFF, 0xF8, 0xE1),
             line_color=OR, line_width_pt=1.5)
    add_rect(slide, Inches(10.6), Inches(1.5), Inches(2.4), Inches(0.45), OR)
    add_text_box(slide, "BASE DE DONNÉES", Inches(10.6), Inches(1.5), Inches(2.4), Inches(0.5),
                 font_size=10, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    items_db = ["PostgreSQL 16", "10 tables", "Hébergé en local\n(port 5433)", "Schéma :\ngestion_palmeraie"]
    bullet_box(slide, items_db, Inches(10.75), Inches(2.1), Inches(2.1), Inches(3.5),
               font_size=11, bullet="▸")
    add_text_box(slide, "Données persistantes", Inches(10.65), Inches(5.8),
                 Inches(2.2), Inches(0.6), font_size=10, italic=True, color=GRIS_TEXTE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 6 – PROCESSUS DE DÉVELOPPEMENT
# ═══════════════════════════════════════════════════════════════════
def slide_06_processus():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Processus de Développement",
                 "Méthodologie itérative adoptée")
    slide_footer(slide, 6)

    etapes = [
        ("1", "Analyse\ndes Besoins",
         "• Recueil des exigences\n• Identification des acteurs\n• Définition du périmètre fonctionnel"),
        ("2", "Conception\nde la BD",
         "• Modèle Conceptuel (MCD)\n• Schéma relationnel (10 tables)\n• Définition des relations"),
        ("3", "Développement\nBackend",
         "• Création des apps Django\n• Modèles → Migrations\n• Sérialisation (DRF)\n• Endpoints REST"),
        ("4", "Développement\nFrontend",
         "• Composants React\n• Connexion API (Axios)\n• Gestion des états (Context)"),
        ("5", "Tests &\nIntégration",
         "• Tests des endpoints API\n• Vérification des formulaires\n• Cohérence des données"),
        ("6", "Déploiement\n& Livraison",
         "• Configuration finale\n• Documentation\n• Présentation"),
    ]

    colors_etapes = [BLEU_FONCE, VERT_FONCE, BLEU_FONCE, VERT_FONCE, BLEU_FONCE, VERT_FONCE]
    card_w = Inches(1.95)
    card_h = Inches(4.6)
    gap = Inches(0.12)

    for idx, (num, titre, detail) in enumerate(etapes):
        x = Inches(0.3) + idx * (card_w + gap)
        color = colors_etapes[idx]

        add_rect(slide, x, Inches(1.45), card_w, card_h, BLANC, line_color=color, line_width_pt=1.2)
        add_rect(slide, x, Inches(1.45), card_w, Inches(0.55), color)

        # Numéro
        add_text_box(slide, f"Étape {num}", x, Inches(1.45), card_w, Inches(0.55),
                     font_size=12, bold=True, color=OR, align=PP_ALIGN.CENTER)

        # Titre étape
        add_text_box(slide, titre, x, Inches(2.1), card_w, Inches(0.75),
                     font_size=11, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Détails
        add_text_box(slide, detail, x + Inches(0.1), Inches(2.95), card_w - Inches(0.2), Inches(2.8),
                     font_size=9.5, color=GRIS_TEXTE, align=PP_ALIGN.LEFT)

    # Flèche de progression
    add_rect(slide, Inches(0.55), Inches(6.28), Inches(12.2), Inches(0.06), OR)
    add_text_box(slide, "PROGRESSION  ➜", Inches(0.3), Inches(6.38), Inches(3), Inches(0.3),
                 font_size=9, bold=True, color=OR)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 7 – OUTILS ET TECHNOLOGIES
# ═══════════════════════════════════════════════════════════════════
def slide_07_outils():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Outils et Technologies Utilisés",
                 "Stack technique complet de l'application")
    slide_footer(slide, 7)

    categories = [
        (BLEU_FONCE, "⚙️  Backend", [
            ("Django 5.2", "Framework Python – gestion des apps, ORM, routing"),
            ("Django REST Framework 3.16", "Construction des API REST + sérialisation JSON"),
            ("Token Authentication", "Sécurisation des endpoints par jeton d'accès"),
            ("django-cors-headers", "Gestion des requêtes cross-origin (CORS)"),
            ("psycopg2-binary", "Pilote de connexion Python ↔ PostgreSQL"),
        ]),
        (VERT_FONCE, "🖥️  Frontend", [
            ("React 19.2", "Bibliothèque UI – composants réactifs"),
            ("Vite 7.3", "Outil de build ultra-rapide pour le développement"),
            ("React Router DOM 7.13", "Navigation entre les pages de l'application"),
            ("Axios", "Client HTTP pour consommer l'API REST"),
            ("ESLint 9", "Analyse statique et qualité du code JavaScript"),
        ]),
        (RGBColor(0x6A, 0x1A, 0x4C), "🗄️  Base de données & Outils", [
            ("PostgreSQL 16", "SGBD relationnel – stockage de toutes les données"),
            ("pgAdmin 4", "Interface graphique d'administration de la BD"),
            ("VS Code", "Éditeur de code principal (backend + frontend)"),
            ("Git / GitHub", "Contrôle de version et collaboration"),
            ("Postman", "Test et débogage des endpoints API"),
        ]),
    ]

    for col_idx, (color, titre, items) in enumerate(categories):
        x = Inches(0.3) + col_idx * Inches(4.35)
        # En-tête catégorie
        add_rect(slide, x, Inches(1.45), Inches(4.15), Inches(0.5), color)
        add_text_box(slide, titre, x, Inches(1.45), Inches(4.15), Inches(0.5),
                     font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

        for row_idx, (outil, desc) in enumerate(items):
            y = Inches(2.05) + row_idx * Inches(0.96)
            bg = VERT_PALE if row_idx % 2 == 0 else BLANC
            add_rect(slide, x, y, Inches(4.15), Inches(0.9), bg, line_color=color, line_width_pt=0.5)
            add_text_box(slide, outil, x + Inches(0.12), y + Inches(0.05),
                         Inches(3.9), Inches(0.35), font_size=11, bold=True, color=color)
            add_text_box(slide, desc, x + Inches(0.12), y + Inches(0.42),
                         Inches(3.9), Inches(0.42), font_size=9.5, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 8 – MODULES DE L'APPLICATION
# ═══════════════════════════════════════════════════════════════════
def slide_08_modules():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Modules de l'Application",
                 "Six modules fonctionnels interconnectés")
    slide_footer(slide, 8)

    modules = [
        ("🌴", "Récoltes", BLEU_FONCE,
         "Module central de l'application.\nGestion des fiches de récolte journalières avec superviseurs, récolteurs, types de régimes et récépissés de vente.",
         ["Créer / modifier une fiche de récolte", "Associer des récolteurs par secteur", "Saisir les régimes (grands/moyens/petits)", "Enregistrer les ventes (client, poids, montant)", "Exporter les données"]),
        ("👷", "Travaux", VERT_FONCE,
         "Suivi des activités agricoles sur la palmeraie.\nGestion des équipes, des consommables et des tâches effectuées par secteur.",
         ["Créer une fiche de travaux", "Lister les consommables utilisés", "Répartir les tâches par personne", "Associer des secteurs couverts", "Export et consultation"]),
        ("🔧", "Matériels", RGBColor(0x4A, 0x14, 0x8C),
         "Inventaire des équipements de la palmeraie.\nSuivi de l'état physique et du statut d'utilisation de chaque matériel.",
         ["Ajouter / modifier un équipement", "Suivre l'état (bon / mauvais / ...)", "Suivre le statut (en service / ...)", "Numérotation automatique", "Liste complète du parc"]),
        ("👥", "Récolteurs", RGBColor(0x1A, 0x23, 0x7E),
         "Annuaire des récolteurs employés.\nCodes automatiques (REC-001...) et historique de leurs performances de récolte.",
         ["Enregistrer un récolteur", "Code généré automatiquement", "Lieu de résidence", "Statistiques de récolte / an", "Export de la liste"]),
        ("🗺️", "Secteurs", RGBColor(0x00, 0x60, 0x64),
         "Gestion des zones géographiques de la palmeraie.\nSuperficie, type de sol, situation de relief avec codes automatiques.",
         ["Créer / modifier un secteur", "Superficie en hectares", "Type de sol et relief", "Code généré automatiquement", "Export des données"]),
        ("📊", "Tableau de Bord", RGBColor(0xBF, 0x36, 0x0C),
         "Vue analytique et décisionnelle.\nGraphiques et KPIs filtrables par année, secteur, récolteur et type de régime.",
         ["Total régimes récoltés", "Chiffre d'affaires global", "Production par secteur", "Graphique mensuel", "Filtres dynamiques"]),
    ]

    cols = [modules[:3], modules[3:]]
    for col_idx, col in enumerate(cols):
        x_base = Inches(0.25) + col_idx * Inches(6.55)
        for row_idx, (icon, nom, color, desc, feats) in enumerate(col):
            y = Inches(1.45) + row_idx * Inches(1.95)
            add_rect(slide, x_base, y, Inches(6.35), Inches(1.85), BLANC,
                     line_color=color, line_width_pt=1.2)
            # Bandeau couleur
            add_rect(slide, x_base, y, Inches(0.6), Inches(1.85), color)
            add_text_box(slide, icon, x_base, y + Inches(0.6), Inches(0.6), Inches(0.65),
                         font_size=18, align=PP_ALIGN.CENTER)
            # Nom du module
            add_text_box(slide, nom, x_base + Inches(0.7), y + Inches(0.05),
                         Inches(5.5), Inches(0.38), font_size=14, bold=True, color=color)
            # Description courte
            add_text_box(slide, desc.split('\n')[0],
                         x_base + Inches(0.7), y + Inches(0.45), Inches(5.5), Inches(1.3),
                         font_size=9.5, color=GRIS_TEXTE)
            # Features
            feat_text = "  •  ".join(feats[:3])
            add_text_box(slide, feat_text, x_base + Inches(0.7), y + Inches(1.2),
                         Inches(5.5), Inches(0.55), font_size=8.5, italic=True, color=color)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 9 – SCHÉMA ENTITÉ-RELATION
# ═══════════════════════════════════════════════════════════════════
def slide_09_er():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x0D, 0x19, 0x2B))
    slide_header(slide, "Schéma de la Base de Données",
                 "Relations entre les 10 tables")
    slide_footer(slide, 9)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.22), BLEU_FONCE)
    add_rect(slide, 0, Inches(1.22), SLIDE_W, Inches(0.065), OR)

    def er_table(x, y, w, name, color, fields):
        add_rect(slide, x, y, w, Inches(0.38), color)
        add_text_box(slide, name, x, y, w, Inches(0.38),
                     font_size=10, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        row_h = Inches(0.27)
        for i, fld in enumerate(fields):
            bg = RGBColor(0x1E, 0x2E, 0x42) if i % 2 == 0 else RGBColor(0x16, 0x23, 0x34)
            add_rect(slide, x, y + Inches(0.38) + i * row_h, w, row_h, bg)
            add_text_box(slide, fld, x + Inches(0.08), y + Inches(0.38) + i * row_h,
                         w - Inches(0.1), row_h, font_size=8, color=RGBColor(0xBB, 0xDE, 0xFF))
        return y + Inches(0.38) + len(fields) * row_h

    # Tables
    er_table(Inches(0.15), Inches(1.45), Inches(2.05), "Secteur", VERT_FONCE,
             ["🔑 id", "code (unique)", "nom", "superficie_ha", "type_sol", "situation_relief"])

    er_table(Inches(0.15), Inches(4.05), Inches(2.05), "Recolteur", BLEU_FONCE,
             ["🔑 id", "code (REC-xxx)", "nom", "lieu_residence"])

    er_table(Inches(2.45), Inches(1.45), Inches(2.35), "FicheRecolte", RGBColor(0x6A, 0x1A, 0x4C),
             ["🔑 id", "date", "superviseur_general", "bareme_grands (60)", "bareme_moyens (50)",
              "bareme_petits (25)", "depense_nourriture", "depense_transport"])

    er_table(Inches(2.45), Inches(5.1), Inches(2.35), "SuperviseurAdjoint", RGBColor(0x4A, 0x14, 0x8C),
             ["🔑 id", "🔗 fiche_id", "nom", "secteur_ou_recolteur"])

    er_table(Inches(5.05), Inches(1.45), Inches(2.35), "FicheRecolteLigne", RGBColor(0x1A, 0x23, 0x7E),
             ["🔑 id", "🔗 fiche_id", "🔗 recolteur_id (null ok)", "recolteur_nom", "regime_type"])

    er_table(Inches(5.05), Inches(4.45), Inches(2.35), "FicheRecolteDetail", RGBColor(0x00, 0x60, 0x64),
             ["🔑 id", "🔗 ligne_id", "🔗 secteur_id (null ok)", "secteur_code", "quantite"])

    er_table(Inches(7.65), Inches(1.45), Inches(2.35), "FicheRecuVente", RGBColor(0xBF, 0x36, 0x0C),
             ["🔑 id", "🔗 fiche_id", "date", "client", "pesee_kg", "non_conformes_pct", "montant"])

    er_table(Inches(10.25), Inches(1.45), Inches(2.8), "FicheTravaux", VERT_FONCE,
             ["🔑 id", "superviseur_travaux", "nature_travaux", "superficie_couverte_ha",
              "nb_personnes", "🔗🔗 secteurs (M2M)"])

    er_table(Inches(10.25), Inches(4.05), Inches(2.8), "ConsommableTravaux", BLEU_FONCE,
             ["🔑 id", "🔗 fiche_id", "designation", "quantite", "unite", "prix_unitaire"])

    er_table(Inches(10.25), Inches(5.8), Inches(2.8), "MaterielEquipement", RGBColor(0x4A, 0x14, 0x8C),
             ["🔑 id", "numero (unique)", "designation", "quantite", "etat_physique", "statut_utilisation"])

    # Relations (flèches texte)
    arrow_color = OR
    arrows = [
        (Inches(4.8), Inches(2.6), "FK"),
        (Inches(4.8), Inches(5.5), "FK"),
        (Inches(7.4), Inches(2.8), "FK"),
        (Inches(2.45+2.35/2), Inches(5.0), "FK"),
        (Inches(5.05+2.35/2), Inches(4.35), "FK"),
    ]

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 10 – TABLE SECTEUR & RECOLTEUR
# ═══════════════════════════════════════════════════════════════════
def slide_10_tables_sr():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Tables : Secteur & Récolteur",
                 "Entités de base — zones géographiques et personnes")
    slide_footer(slide, 10)

    # Table Secteur
    add_rect(slide, Inches(0.3), Inches(1.45), Inches(5.9), Inches(0.45), VERT_FONCE)
    add_text_box(slide, "TABLE : Secteur", Inches(0.3), Inches(1.45), Inches(5.9), Inches(0.45),
                 font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_s = ["Attribut", "Type", "Description"]
    rows_s = [
        ["id", "INTEGER (PK)", "Clé primaire automatique"],
        ["code", "VARCHAR(20) UNIQUE", "Code auto-généré (ex: GP_1, SEC_2)"],
        ["nom", "VARCHAR(120)", "Nom du secteur"],
        ["superficie_ha", "DECIMAL(8,2)", "Superficie en hectares"],
        ["situation_relief", "VARCHAR(120)", "Description du relief (plaine, colline…)"],
        ["type_sol", "VARCHAR(200)", "Type de sol (argileux, sableux…)"],
        ["created_at", "TIMESTAMP", "Date de création (auto)"],
    ]
    add_table(slide, headers_s, rows_s,
              Inches(0.3), Inches(1.95), Inches(5.9), Inches(2.5),
              header_fill=VERT_FONCE, font_size=9)

    add_text_box(slide,
                 "Relations : ← FicheRecolteDetail (via secteur_id) | ← FicheTravaux (via M2M secteurs_couverts)",
                 Inches(0.3), Inches(4.55), Inches(5.9), Inches(0.4),
                 font_size=8.5, italic=True, color=VERT_FONCE)

    # Table Recolteur
    add_rect(slide, Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.45), BLEU_FONCE)
    add_text_box(slide, "TABLE : Recolteur", Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.45),
                 font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_r = ["Attribut", "Type", "Description"]
    rows_r = [
        ["id", "INTEGER (PK)", "Clé primaire automatique"],
        ["code", "VARCHAR(20) UNIQUE", "Code auto-généré (REC-001, REC-042…)"],
        ["nom", "VARCHAR(120)", "Nom complet du récolteur"],
        ["lieu_residence", "VARCHAR(120)", "Lieu de résidence"],
        ["created_at", "TIMESTAMP", "Date de création (auto)"],
    ]
    add_table(slide, headers_r, rows_r,
              Inches(7.0), Inches(1.95), Inches(5.9), Inches(1.85),
              header_fill=BLEU_FONCE, font_size=9)

    add_text_box(slide,
                 "Relations : ← FicheRecolteLigne (via recolteur_id, nullable)\nCode = 'REC-' + str(pk).zfill(3) (généré automatiquement à la sauvegarde)",
                 Inches(7.0), Inches(3.88), Inches(5.9), Inches(0.65),
                 font_size=8.5, italic=True, color=BLEU_FONCE)

    # Séparateur vertical
    add_rect(slide, Inches(6.55), Inches(1.45), Inches(0.06), Inches(5.6), VERT_CLAIR)

    # Encadré modification
    add_rect(slide, Inches(0.3), Inches(5.05), Inches(12.6), Inches(1.65), BLEU_CIEL,
             line_color=BLEU_FONCE, line_width_pt=1)
    add_text_box(slide, "💡 Attributs à discuter — Ajouts possibles :",
                 Inches(0.45), Inches(5.12), Inches(12.2), Inches(0.38),
                 font_size=11, bold=True, color=BLEU_FONCE)
    add_text_box(slide,
                 "Secteur : age_moyen_plants (années), nb_palmiers, rendement_cible_t_ha, coordonnees_GPS, statut (actif/inactif)\n"
                 "Récolteur : date_naissance, telephone, photo, contrat (journalier/permanent), date_embauche",
                 Inches(0.45), Inches(5.55), Inches(12.2), Inches(1.05),
                 font_size=10, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 11 – TABLES RÉCOLTES (1)
# ═══════════════════════════════════════════════════════════════════
def slide_11_recoltes1():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Tables Récoltes (1/2) : FicheRecolte & SuperviseurAdjoint",
                 "En-tête de la fiche de récolte et ses superviseurs adjoints")
    slide_footer(slide, 11)

    # FicheRecolte
    add_rect(slide, Inches(0.3), Inches(1.45), Inches(7.8), Inches(0.45), RGBColor(0x6A, 0x1A, 0x4C))
    add_text_box(slide, "TABLE : FicheRecolte", Inches(0.3), Inches(1.45), Inches(7.8), Inches(0.45),
                 font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_fr = ["Attribut", "Type", "Valeur par défaut", "Description"]
    rows_fr = [
        ["id", "INTEGER (PK)", "—", "Clé primaire automatique"],
        ["date", "DATE", "—", "Date de la récolte"],
        ["superviseur_general", "VARCHAR(120)", "''", "Nom du superviseur principal"],
        ["bareme_grands", "INT UNSIGNED", "60", "Barème unitaire régimes grands (FCFA/unité)"],
        ["bareme_moyens", "INT UNSIGNED", "50", "Barème unitaire régimes moyens (FCFA/unité)"],
        ["bareme_petits", "INT UNSIGNED", "25", "Barème unitaire régimes petits (FCFA/unité)"],
        ["depense_nourriture", "DECIMAL(12,2)", "0.00", "Dépenses nourriture (FCFA)"],
        ["depense_transport", "DECIMAL(12,2)", "0.00", "Dépenses transport (FCFA)"],
        ["observations", "TEXT", "''", "Notes libres du superviseur"],
        ["created_at", "TIMESTAMP", "now()", "Date de création automatique"],
    ]
    add_table(slide, headers_fr, rows_fr,
              Inches(0.3), Inches(1.95), Inches(7.8), Inches(3.4),
              header_fill=RGBColor(0x6A, 0x1A, 0x4C), font_size=8.5)

    # SuperviseurAdjoint
    add_rect(slide, Inches(8.3), Inches(1.45), Inches(4.7), Inches(0.45), RGBColor(0x4A, 0x14, 0x8C))
    add_text_box(slide, "TABLE : SuperviseurAdjoint", Inches(8.3), Inches(1.45), Inches(4.7), Inches(0.45),
                 font_size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_sa = ["Attribut", "Type", "Description"]
    rows_sa = [
        ["id", "INTEGER (PK)", "Clé primaire"],
        ["fiche_id", "FK → FicheRecolte", "Fiche parente (CASCADE)"],
        ["nom", "VARCHAR(120)", "Nom du superviseur adjoint"],
        ["secteur_ou_recolteur", "VARCHAR(120)", "Secteur ou récolteur supervisé"],
    ]
    add_table(slide, headers_sa, rows_sa,
              Inches(8.3), Inches(1.95), Inches(4.7), Inches(1.7),
              header_fill=RGBColor(0x4A, 0x14, 0x8C), font_size=9)

    # Encadré modification
    add_rect(slide, Inches(0.3), Inches(5.5), Inches(12.6), Inches(1.65), BLEU_CIEL,
             line_color=BLEU_FONCE, line_width_pt=1)
    add_text_box(slide, "💡 Attributs à discuter — Ajouts possibles :",
                 Inches(0.45), Inches(5.57), Inches(12.2), Inches(0.38),
                 font_size=11, bold=True, color=BLEU_FONCE)
    add_text_box(slide,
                 "FicheRecolte : heure_debut / heure_fin, conditions_meteo, nb_palmiers_recoltes, surface_recoltee_ha, statut (brouillon/validé)\n"
                 "SuperviseurAdjoint : matricule, telephone",
                 Inches(0.45), Inches(6.0), Inches(12.2), Inches(1.0),
                 font_size=10, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 12 – TABLES RÉCOLTES (2)
# ═══════════════════════════════════════════════════════════════════
def slide_12_recoltes2():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Tables Récoltes (2/2) : Ligne · Détail · RecuVente",
                 "Lignes par récolteur, détails par secteur et récépissés de vente")
    slide_footer(slide, 12)

    # FicheRecolteLigne
    add_rect(slide, Inches(0.3), Inches(1.45), Inches(4.1), Inches(0.45), RGBColor(0x1A, 0x23, 0x7E))
    add_text_box(slide, "TABLE : FicheRecolteLigne", Inches(0.3), Inches(1.45), Inches(4.1), Inches(0.45),
                 font_size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_rl = ["Attribut", "Type", "Description"]
    rows_rl = [
        ["id", "INTEGER (PK)", "Clé primaire"],
        ["fiche_id", "FK → FicheRecolte", "Fiche parente (CASCADE)"],
        ["recolteur_id", "FK → Recolteur (null)", "Récolteur (peut être null)"],
        ["recolteur_nom", "VARCHAR(120)", "Nom (copié si récolteur null)"],
        ["regime_type", "VARCHAR(10)", "grands | moyens | petits"],
    ]
    add_table(slide, headers_rl, rows_rl,
              Inches(0.3), Inches(1.95), Inches(4.1), Inches(1.9),
              header_fill=RGBColor(0x1A, 0x23, 0x7E), font_size=8.5)

    # FicheRecolteDetail
    add_rect(slide, Inches(4.65), Inches(1.45), Inches(4.1), Inches(0.45), RGBColor(0x00, 0x60, 0x64))
    add_text_box(slide, "TABLE : FicheRecolteDetail", Inches(4.65), Inches(1.45), Inches(4.1), Inches(0.45),
                 font_size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_rd = ["Attribut", "Type", "Description"]
    rows_rd = [
        ["id", "INTEGER (PK)", "Clé primaire"],
        ["ligne_id", "FK → FicheRecolteLigne", "Ligne parente (CASCADE)"],
        ["secteur_id", "FK → Secteur (null)", "Secteur concerné"],
        ["secteur_code", "VARCHAR(20)", "Code copié (dénormalisation)"],
        ["quantite", "INT UNSIGNED", "Nombre de régimes récoltés"],
    ]
    add_table(slide, headers_rd, rows_rd,
              Inches(4.65), Inches(1.95), Inches(4.1), Inches(1.9),
              header_fill=RGBColor(0x00, 0x60, 0x64), font_size=8.5)

    # FicheRecuVente
    add_rect(slide, Inches(9.0), Inches(1.45), Inches(4.0), Inches(0.45), RGBColor(0xBF, 0x36, 0x0C))
    add_text_box(slide, "TABLE : FicheRecuVente", Inches(9.0), Inches(1.45), Inches(4.0), Inches(0.45),
                 font_size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_rv = ["Attribut", "Type", "Description"]
    rows_rv = [
        ["id", "INTEGER (PK)", "Clé primaire"],
        ["fiche_id", "FK → FicheRecolte", "Fiche parente (CASCADE)"],
        ["date", "DATE (null ok)", "Date de la vente"],
        ["client", "VARCHAR(120)", "Nom du client acheteur"],
        ["pesee_kg", "DECIMAL(10,2)", "Poids pesé en kg"],
        ["non_conformes_pct", "DECIMAL(5,2)", "% de régimes non conformes"],
        ["montant", "DECIMAL(12,2)", "Montant encaissé (FCFA)"],
    ]
    add_table(slide, headers_rv, rows_rv,
              Inches(9.0), Inches(1.95), Inches(4.0), Inches(2.3),
              header_fill=RGBColor(0xBF, 0x36, 0x0C), font_size=8.5)

    # Flux de données
    add_rect(slide, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.04), OR)
    add_text_box(slide, "Flux des données :",
                 Inches(0.3), Inches(4.1), Inches(2.5), Inches(0.4),
                 font_size=11, bold=True, color=BLEU_FONCE)
    add_text_box(slide,
                 "FicheRecolte  →  (1 à N)  FicheRecolteLigne  →  (1 à N)  FicheRecolteDetail  +  FicheRecuVente\n"
                 "Une fiche contient plusieurs lignes (une par récolteur × régime). Chaque ligne peut couvrir plusieurs secteurs (détails).",
                 Inches(2.85), Inches(4.1), Inches(10.1), Inches(0.75),
                 font_size=10, color=GRIS_TEXTE)

    # Encadré ajouts possibles
    add_rect(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(1.65), BLEU_CIEL,
             line_color=BLEU_FONCE, line_width_pt=1)
    add_text_box(slide, "💡 Attributs à discuter — Ajouts possibles :",
                 Inches(0.45), Inches(5.07), Inches(12.2), Inches(0.38),
                 font_size=11, bold=True, color=BLEU_FONCE)
    add_text_box(slide,
                 "FicheRecolteLigne : salaire_calcule (montant calculé automatiquement), prime_qualite, nb_heures_travail\n"
                 "FicheRecolteDetail : coordonnees_GPS_palmier, qualite_regime (A/B/C)\n"
                 "FicheRecuVente : reference_facture, mode_paiement (espèce/virement), vehicule_transport",
                 Inches(0.45), Inches(5.5), Inches(12.2), Inches(1.1),
                 font_size=10, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 13 – TABLES TRAVAUX
# ═══════════════════════════════════════════════════════════════════
def slide_13_travaux():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Tables Travaux : FicheTravaux · Consommables · RépartitionTâche",
                 "Gestion des activités agricoles et de la main-d'œuvre")
    slide_footer(slide, 13)

    # FicheTravaux
    add_rect(slide, Inches(0.3), Inches(1.45), Inches(4.0), Inches(0.45), VERT_FONCE)
    add_text_box(slide, "TABLE : FicheTravaux", Inches(0.3), Inches(1.45), Inches(4.0), Inches(0.45),
                 font_size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_ft = ["Attribut", "Type", "Description"]
    rows_ft = [
        ["id", "INTEGER (PK)", "Clé primaire"],
        ["superviseur_travaux", "VARCHAR(120)", "Responsable des travaux"],
        ["nature_travaux", "VARCHAR(255)", "Type de travaux (taille, fertilisation…)"],
        ["superficie_couverte_ha", "DECIMAL(8,2) null", "Surface concernée en ha"],
        ["periode_travaux", "VARCHAR(120)", "Période (ex: Janvier 2025)"],
        ["nb_personnes", "INT UNSIGNED null", "Nombre de personnes mobilisées"],
        ["secteurs_couverts", "M2M → Secteur", "Secteurs concernés (plusieurs)"],
        ["observations", "TEXT", "Notes libres"],
        ["created_at", "TIMESTAMP", "Date de création (auto)"],
    ]
    add_table(slide, headers_ft, rows_ft,
              Inches(0.3), Inches(1.95), Inches(4.0), Inches(3.0),
              header_fill=VERT_FONCE, font_size=8.5)

    # ConsommableTravaux
    add_rect(slide, Inches(4.55), Inches(1.45), Inches(4.0), Inches(0.45), BLEU_FONCE)
    add_text_box(slide, "TABLE : ConsommableTravaux", Inches(4.55), Inches(1.45), Inches(4.0), Inches(0.45),
                 font_size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_ct = ["Attribut", "Type", "Description"]
    rows_ct = [
        ["id", "INTEGER (PK)", "Clé primaire"],
        ["fiche_id", "FK → FicheTravaux", "Fiche parente (CASCADE)"],
        ["designation", "VARCHAR(200)", "Nom du produit/matériel consommé"],
        ["quantite", "DECIMAL(10,2)", "Quantité utilisée"],
        ["unite", "VARCHAR(40)", "Unité (kg, L, sac, …)"],
        ["prix_unitaire", "DECIMAL(12,2)", "Prix par unité (FCFA)"],
        ["created_at", "TIMESTAMP", "Date de création"],
    ]
    add_table(slide, headers_ct, rows_ct,
              Inches(4.55), Inches(1.95), Inches(4.0), Inches(2.6),
              header_fill=BLEU_FONCE, font_size=8.5)

    # RepartitionTache
    add_rect(slide, Inches(8.8), Inches(1.45), Inches(4.2), Inches(0.45), RGBColor(0x4A, 0x14, 0x8C))
    add_text_box(slide, "TABLE : RepartitionTache", Inches(8.8), Inches(1.45), Inches(4.2), Inches(0.45),
                 font_size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_rt = ["Attribut", "Type", "Description"]
    rows_rt = [
        ["id", "INTEGER (PK)", "Clé primaire"],
        ["fiche_id", "FK → FicheTravaux", "Fiche parente (CASCADE)"],
        ["nom_prenom", "VARCHAR(120)", "Nom de l'ouvrier"],
        ["nature_taches", "VARCHAR(200)", "Tâche effectuée"],
        ["quantite", "DECIMAL(10,2)", "Quantité de travail réalisée"],
        ["prix_unitaire", "DECIMAL(12,2)", "Tarif unitaire (FCFA)"],
        ["created_at", "TIMESTAMP", "Date de création"],
    ]
    add_table(slide, headers_rt, rows_rt,
              Inches(8.8), Inches(1.95), Inches(4.2), Inches(2.6),
              header_fill=RGBColor(0x4A, 0x14, 0x8C), font_size=8.5)

    # Encadré ajouts
    add_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.65), BLEU_CIEL,
             line_color=BLEU_FONCE, line_width_pt=1)
    add_text_box(slide, "💡 Attributs à discuter — Ajouts possibles :",
                 Inches(0.45), Inches(5.17), Inches(12.2), Inches(0.38),
                 font_size=11, bold=True, color=BLEU_FONCE)
    add_text_box(slide,
                 "FicheTravaux : date_debut, date_fin, cout_total_calcule (auto), statut (planifié/en cours/terminé), type_travaux (ENUM)\n"
                 "ConsommableTravaux : fournisseur, numero_lot, date_peremption\n"
                 "RepartitionTache : salaire_total_calcule (quantite × prix_unitaire), matricule_ouvrier",
                 Inches(0.45), Inches(5.6), Inches(12.2), Inches(1.05),
                 font_size=10, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 14 – TABLE MATÉRIEL
# ═══════════════════════════════════════════════════════════════════
def slide_14_materiel():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Table : MaterielEquipement",
                 "Inventaire du matériel et des équipements de la palmeraie")
    slide_footer(slide, 14)

    add_rect(slide, Inches(2.0), Inches(1.45), Inches(9.3), Inches(0.45), RGBColor(0x4A, 0x14, 0x8C))
    add_text_box(slide, "TABLE : MaterielEquipement", Inches(2.0), Inches(1.45), Inches(9.3), Inches(0.45),
                 font_size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    headers_m = ["Attribut", "Type Django", "Type SQL", "Contraintes", "Description"]
    rows_m = [
        ["id", "AutoField", "INTEGER", "PK, AUTO", "Clé primaire automatique"],
        ["numero", "PositiveIntegerField", "INT UNSIGNED", "UNIQUE, NOT NULL", "Numéro d'inventaire unique"],
        ["designation", "CharField(200)", "VARCHAR(200)", "blank=True", "Nom / description du matériel"],
        ["quantite", "PositiveIntegerField", "INT UNSIGNED", "default=0", "Quantité en stock"],
        ["etat_physique", "CharField(120)", "VARCHAR(120)", "blank=True", "État : bon, mauvais, à réparer…"],
        ["statut_utilisation", "CharField(120)", "VARCHAR(120)", "blank=True", "Statut : en service, en panne…"],
        ["created_at", "DateTimeField", "TIMESTAMP", "auto_now_add=True", "Date d'ajout dans le système"],
    ]
    add_table(slide, headers_m, rows_m,
              Inches(0.4), Inches(2.0), Inches(12.5), Inches(2.8),
              header_fill=RGBColor(0x4A, 0x14, 0x8C), font_size=9.5)

    # Note
    add_rect(slide, Inches(0.4), Inches(4.95), Inches(12.5), Inches(0.06), OR)
    add_text_box(slide, "Table autonome — aucune relation avec les autres tables actuellement.",
                 Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.4),
                 font_size=10, italic=True, color=GRIS_TEXTE)

    # Encadré ajouts
    add_rect(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.65), BLEU_CIEL,
             line_color=BLEU_FONCE, line_width_pt=1)
    add_text_box(slide, "💡 Attributs à discuter — Ajouts possibles :",
                 Inches(0.55), Inches(5.67), Inches(12.0), Inches(0.38),
                 font_size=11, bold=True, color=BLEU_FONCE)
    add_text_box(slide,
                 "date_acquisition, valeur_achat (FCFA), fournisseur, date_derniere_maintenance, date_prochaine_maintenance\n"
                 "categorie (véhicule / outil / équipement lourd / petit matériel), localisation (secteur ou dépôt), responsable\n"
                 "Relation possible : lier le matériel aux FicheTravaux pour tracer l'utilisation par chantier",
                 Inches(0.55), Inches(6.1), Inches(12.0), Inches(1.05),
                 font_size=10, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 15 – FONCTIONNALITÉS PRINCIPALES
# ═══════════════════════════════════════════════════════════════════
def slide_15_fonctionnalites():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Fonctionnalités Principales",
                 "Ce que l'utilisateur peut faire dans l'application")
    slide_footer(slide, 15)

    features = [
        (BLEU_FONCE, "Authentification sécurisée",
         ["Connexion par identifiant / mot de passe",
          "Jeton JWT d'accès (Token Auth DRF)",
          "Toutes les pages protégées",
          "Déconnexion sécurisée"]),
        (VERT_FONCE, "CRUD complet sur toutes les entités",
         ["Créer, consulter, modifier, supprimer",
          "Formulaires de saisie validés côté client",
          "Messages de confirmation et d'erreur",
          "Interface responsive"]),
        (RGBColor(0x6A, 0x1A, 0x4C), "Fiches structurées",
         ["Saisie de fiches multi-niveaux (entête → lignes → détails)",
          "Plusieurs superviseurs adjoints par fiche",
          "Plusieurs secteurs par ligne de récolte",
          "Récépissés de vente attachés"]),
        (RGBColor(0x1A, 0x23, 0x7E), "Exports de données",
         ["Export disponible pour : récoltes, travaux,",
          "  secteurs, récolteurs",
          "Format téléchargeable (CSV/Excel)",
          "Filtres par année ou secteur"]),
        (RGBColor(0x00, 0x60, 0x64), "Codes automatiques",
         ["REC-001 … REC-999 pour les récolteurs",
          "GP_1, SEC_2 … pour les secteurs",
          "Évite les doublons et saisies manuelles",
          "Immutable après création"]),
        (RGBColor(0xBF, 0x36, 0x0C), "Tableau de bord analytique",
         ["KPIs clés (total régimes, CA, nb récolteurs actifs)",
          "Graphique de production mensuelle",
          "Filtres : année, secteur, récolteur, régime",
          "Actualisation en temps réel"]),
    ]

    cols = [features[:3], features[3:]]
    for col_idx, col in enumerate(cols):
        x_base = Inches(0.25) + col_idx * Inches(6.55)
        for row_idx, (color, titre, items) in enumerate(col):
            y = Inches(1.45) + row_idx * Inches(1.88)
            add_rect(slide, x_base, y, Inches(6.35), Inches(1.78), BLANC,
                     line_color=color, line_width_pt=1.2)
            add_rect(slide, x_base, y, Inches(6.35), Inches(0.4), color)
            add_text_box(slide, titre, x_base + Inches(0.1), y + Inches(0.04),
                         Inches(6.1), Inches(0.35), font_size=12, bold=True, color=BLANC)
            for i, item in enumerate(items):
                add_text_box(slide, f"✓  {item}",
                             x_base + Inches(0.1), y + Inches(0.48) + i * Inches(0.3),
                             Inches(6.1), Inches(0.28), font_size=10, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 16 – TABLEAU DE BORD
# ═══════════════════════════════════════════════════════════════════
def slide_16_dashboard():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Tableau de Bord & Analyses",
                 "Indicateurs clés et visualisations pour la prise de décision")
    slide_footer(slide, 16)

    # KPIs
    kpis = [
        (BLEU_FONCE, "🌴", "Total Régimes", "Somme de tous les régimes récoltés (grands + moyens + petits)"),
        (VERT_FONCE, "💰", "Chiffre d'Affaires", "Total des montants encaissés sur les récépissés de vente"),
        (RGBColor(0xBF, 0x36, 0x0C), "👷", "Récolteurs Actifs", "Nombre de récolteurs ayant participé à au moins une récolte"),
        (RGBColor(0x6A, 0x1A, 0x4C), "📅", "Fiches Récolte", "Nombre total de fiches de récolte enregistrées"),
    ]
    for idx, (color, icon, label, desc) in enumerate(kpis):
        x = Inches(0.3) + idx * Inches(3.2)
        add_rect(slide, x, Inches(1.45), Inches(3.0), Inches(1.5), BLANC,
                 line_color=color, line_width_pt=1.5)
        add_rect(slide, x, Inches(1.45), Inches(0.55), Inches(1.5), color)
        add_text_box(slide, icon, x, Inches(1.65), Inches(0.55), Inches(0.7),
                     font_size=18, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x + Inches(0.65), Inches(1.52),
                     Inches(2.25), Inches(0.42), font_size=12, bold=True, color=color)
        add_text_box(slide, desc, x + Inches(0.65), Inches(1.97),
                     Inches(2.25), Inches(0.85), font_size=8.5, color=GRIS_TEXTE)

    # Section Graphique
    add_text_box(slide, "📊  Graphique de production mensuelle",
                 Inches(0.3), Inches(3.15), Inches(8.5), Inches(0.45),
                 font_size=14, bold=True, color=BLEU_FONCE)
    add_rect(slide, Inches(0.3), Inches(3.65), Inches(8.5), Inches(3.0), BLANC,
             line_color=BLEU_FONCE, line_width_pt=1)

    # Simulation graphique bar chart
    mois = ["Jan", "Fév", "Mar", "Avr", "Mai", "Jun", "Jul", "Aoû", "Sep", "Oct", "Nov", "Déc"]
    heights_rel = [0.3, 0.5, 0.7, 0.9, 0.6, 0.8, 1.0, 0.75, 0.65, 0.85, 0.4, 0.55]
    bar_w = Inches(0.52)
    chart_bottom = Inches(6.45)
    chart_max_h = Inches(2.3)
    chart_left = Inches(0.55)

    for i, (m, h) in enumerate(zip(mois, heights_rel)):
        bar_h = chart_max_h * h
        x = chart_left + i * Inches(0.6)
        y = chart_bottom - bar_h
        color_bar = VERT_CLAIR if i % 2 == 0 else BLEU_FONCE
        add_rect(slide, x, y, bar_w, bar_h, color_bar)
        add_text_box(slide, m, x, chart_bottom + Inches(0.02), bar_w, Inches(0.25),
                     font_size=7, align=PP_ALIGN.CENTER, color=GRIS_TEXTE)

    add_text_box(slide, "Nombre de régimes récoltés par mois (simulation de données)",
                 Inches(0.5), Inches(3.72), Inches(8.0), Inches(0.3),
                 font_size=8, italic=True, color=GRIS_TEXTE)

    # Filtres disponibles
    add_text_box(slide, "⚙️  Filtres disponibles",
                 Inches(9.2), Inches(3.15), Inches(3.8), Inches(0.45),
                 font_size=14, bold=True, color=VERT_FONCE)

    filtres = [
        ("📆", "Année", "Sélectionner une année (ex: 2024, 2025)"),
        ("🗺️", "Secteur", "Filtrer par secteur géographique"),
        ("👷", "Récolteur", "Voir les données d'un récolteur précis"),
        ("🌴", "Type régime", "grands / moyens / petits"),
    ]
    for idx, (icon, label, desc) in enumerate(filtres):
        y = Inches(3.7) + idx * Inches(0.82)
        add_rect(slide, Inches(9.2), y, Inches(3.8), Inches(0.72), BLANC,
                 line_color=VERT_CLAIR, line_width_pt=0.8)
        add_text_box(slide, f"{icon} {label}", Inches(9.3), y + Inches(0.05),
                     Inches(3.5), Inches(0.32), font_size=11, bold=True, color=VERT_FONCE)
        add_text_box(slide, desc, Inches(9.3), y + Inches(0.38),
                     Inches(3.5), Inches(0.28), font_size=9, color=GRIS_TEXTE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 17 – SÉCURITÉ & API
# ═══════════════════════════════════════════════════════════════════
def slide_17_api():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "API REST & Sécurité",
                 "Points d'accès de l'application et mécanismes de protection")
    slide_footer(slide, 17)

    endpoints = [
        ("Auth",    "POST",  "/api/auth/login/",                     "Connexion — obtenir un token d'accès",         BLEU_FONCE),
        ("Secteurs","GET",   "/api/secteurs/",                       "Lister tous les secteurs",                     VERT_FONCE),
        ("Secteurs","POST",  "/api/secteurs/",                       "Créer un nouveau secteur",                     VERT_FONCE),
        ("Secteurs","GET",   "/api/secteurs/{id}/export/?year=YYYY", "Exporter les données d'un secteur",            VERT_FONCE),
        ("Récolteurs","GET", "/api/recolteurs/",                     "Lister les récolteurs",                        BLEU_FONCE),
        ("Récolteurs","GET", "/api/recolteurs/stats/?year=YYYY",     "Statistiques par récolteur et par année",      BLEU_FONCE),
        ("Récoltes","GET",   "/api/recoltes/",                       "Lister les fiches de récolte",                 RGBColor(0x6A, 0x1A, 0x4C)),
        ("Récoltes","POST",  "/api/recoltes/",                       "Créer une fiche de récolte complète",          RGBColor(0x6A, 0x1A, 0x4C)),
        ("Récoltes","GET",   "/api/recoltes/analytics/?year=YYYY",   "Données analytiques de récolte",               RGBColor(0x6A, 0x1A, 0x4C)),
        ("Travaux", "GET",   "/api/travaux/",                        "Lister les fiches de travaux",                 RGBColor(0x1A, 0x23, 0x7E)),
        ("Matériels","GET",  "/api/materiels/",                      "Lister les équipements",                       RGBColor(0x4A, 0x14, 0x8C)),
        ("Dashboard","GET",  "/api/dashboard/summary/",              "Résumé global filtrable (année, secteur, …)",  RGBColor(0xBF, 0x36, 0x0C)),
    ]

    headers_api = ["Module", "Méthode", "Endpoint", "Description"]
    rows_api = [[m, meth, ep, desc] for m, meth, ep, desc, _ in endpoints]

    add_rect(slide, Inches(0.3), Inches(1.45), Inches(12.7), Inches(0.45), BLEU_FONCE)
    add_text_box(slide, "Endpoints disponibles (API REST — authentification Token requise sauf /login/)",
                 Inches(0.3), Inches(1.45), Inches(12.7), Inches(0.45),
                 font_size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_table(slide, headers_api, rows_api,
              Inches(0.3), Inches(1.95), Inches(12.7), Inches(4.1),
              header_fill=BLEU_FONCE, font_size=8.5)

    # Sécurité
    add_rect(slide, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.75), RGBColor(0xFF, 0xF8, 0xE1),
             line_color=OR, line_width_pt=1)
    add_text_box(slide, "🔒 Sécurité :  Token Authentication DRF  ·  IsAuthenticated sur tous les endpoints  ·  CORS configuré pour le frontend local",
                 Inches(0.5), Inches(6.3), Inches(12.2), Inches(0.55),
                 font_size=10, bold=True, color=BLEU_FONCE)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 18 – CONCLUSION
# ═══════════════════════════════════════════════════════════════════
def slide_18_conclusion():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GRIS_CLAIR)
    slide_header(slide, "Conclusion & Perspectives",
                 "Bilan du projet et axes d'amélioration futurs")
    slide_footer(slide, 18)

    # Bilan
    add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.0), Inches(0.45), BLEU_FONCE)
    add_text_box(slide, "✅  Ce qui a été réalisé", Inches(0.3), Inches(1.45), Inches(6.0), Inches(0.45),
                 font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    bilan = [
        "Application web full-stack fonctionnelle (React + Django)",
        "6 modules couvrant l'ensemble des besoins de gestion",
        "10 tables relationnelles avec contraintes d'intégrité",
        "API REST sécurisée avec authentification par token",
        "Tableau de bord analytique avec filtres dynamiques",
        "Fonctionnalités d'export de données",
        "Codes automatiques pour récolteurs et secteurs",
    ]
    bullet_box(slide, bilan, Inches(0.4), Inches(2.0), Inches(5.8), Inches(3.5),
               font_size=11, bullet="✓")

    # Perspectives
    add_rect(slide, Inches(6.9), Inches(1.45), Inches(6.1), Inches(0.45), VERT_FONCE)
    add_text_box(slide, "🚀  Perspectives d'évolution", Inches(6.9), Inches(1.45), Inches(6.1), Inches(0.45),
                 font_size=13, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    perspectives = [
        "Application mobile (React Native / Flutter)",
        "Rapports PDF automatisés par période",
        "Intégration d'un module de paie des récolteurs",
        "Géolocalisation des secteurs (carte interactive)",
        "Module de maintenance du matériel (alertes)",
        "Prévisions de production (IA / Machine Learning)",
        "Déploiement cloud (hébergement en ligne)",
    ]
    bullet_box(slide, perspectives, Inches(7.0), Inches(2.0), Inches(5.9), Inches(3.5),
               font_size=11, bullet="→")

    # Ligne séparatrice
    add_rect(slide, Inches(6.55), Inches(1.45), Inches(0.06), Inches(4.1), OR)

    # Citation finale
    add_rect(slide, Inches(0.3), Inches(5.8), Inches(12.7), Inches(1.3), BLEU_FONCE)
    add_text_box(slide,
                 "\"Ce projet pose les bases d'un système de gestion intelligent de la palmeraie.\n"
                 " Il centralise les données, améliore la traçabilité et offre une aide à la décision\n"
                 " grâce aux analyses en temps réel.\"",
                 Inches(0.6), Inches(5.9), Inches(12.0), Inches(1.1),
                 font_size=12, italic=True, color=BLANC, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 19 – QUESTIONS
# ═══════════════════════════════════════════════════════════════════
def slide_19_questions():
    slide = prs.slides.add_slide(BLANK)

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLEU_FONCE)
    add_rect(slide, 0, Inches(3.5), SLIDE_W, Inches(0.12), OR)

    add_text_box(slide, "🌴", Inches(5.5), Inches(0.4), Inches(2.3), Inches(2.5),
                 font_size=90, align=PP_ALIGN.CENTER, color=OR)

    add_text_box(slide, "Merci pour votre attention",
                 Inches(0.5), Inches(3.75), SLIDE_W - Inches(1.0), Inches(0.8),
                 font_size=30, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Des questions ?",
                 Inches(0.5), Inches(4.6), SLIDE_W - Inches(1.0), Inches(0.7),
                 font_size=24, italic=True, color=OR, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "ABO KOKOUTSE ADZEYI HENOCH\n"
                 "Master 2 BDGL — UFHB · Abidjan\n"
                 "Année académique 2025–2026",
                 Inches(0.5), Inches(5.55), SLIDE_W - Inches(1.0), Inches(1.2),
                 font_size=13, color=RGBColor(0xBB, 0xDE, 0xFF), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), RGBColor(0x0D, 0x1A, 0x2D))
    add_text_box(slide, "Système Intelligent de Gestion du Palmier à Huile",
                 Inches(0.3), Inches(7.16), Inches(13), Inches(0.3),
                 font_size=9, color=OR, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
#  GÉNÉRATION
# ═══════════════════════════════════════════════════════════════════
print("Génération des slides...")
slide_01_titre()
print("  [1/19] ✓ Titre")
slide_02_plan()
print("  [2/19] ✓ Plan")
slide_03_contexte()
print("  [3/19] ✓ Contexte & Problématique")
slide_04_objectifs()
print("  [4/19] ✓ Objectifs")
slide_05_architecture()
print("  [5/19] ✓ Architecture")
slide_06_processus()
print("  [6/19] ✓ Processus de développement")
slide_07_outils()
print("  [7/19] ✓ Outils & Technologies")
slide_08_modules()
print("  [8/19] ✓ Modules")
slide_09_er()
print("  [9/19] ✓ Schéma ER")
slide_10_tables_sr()
print("  [10/19] ✓ Tables Secteur & Récolteur")
slide_11_recoltes1()
print("  [11/19] ✓ Tables Récoltes 1/2")
slide_12_recoltes2()
print("  [12/19] ✓ Tables Récoltes 2/2")
slide_13_travaux()
print("  [13/19] ✓ Tables Travaux")
slide_14_materiel()
print("  [14/19] ✓ Table Matériel")
slide_15_fonctionnalites()
print("  [15/19] ✓ Fonctionnalités")
slide_16_dashboard()
print("  [16/19] ✓ Tableau de bord")
slide_17_api()
print("  [17/19] ✓ API & Sécurité")
slide_18_conclusion()
print("  [18/19] ✓ Conclusion")
slide_19_questions()
print("  [19/19] ✓ Questions")

output_path = r"e:\UFHB\MASTER2_BDGL\Projet_Memoire\theme2\Projet_Gestion_Palmeraie\Presentation_Palmeraie_ABO_KOKOUTSE.pptx"
prs.save(output_path)
print(f"\n✅  Fichier sauvegardé : {output_path}")
