"""
Génération de la présentation PowerPoint — Projet Gestion Palmeraie
Auteur  : généré automatiquement
Usage   : python generate_pptx.py
Sortie  : Presentation_Gestion_Palmeraie.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette de couleurs ──────────────────────────────────────────────────────
VERT_FONCE   = RGBColor(0x1B, 0x5E, 0x20)   # #1B5E20
VERT_MOYEN   = RGBColor(0x2E, 0x7D, 0x32)   # #2E7D32
VERT_CLAIR   = RGBColor(0xA5, 0xD6, 0xA7)   # #A5D6A7
VERT_PALE    = RGBColor(0xE8, 0xF5, 0xE9)   # #E8F5E9
OR           = RGBColor(0xF9, 0xA8, 0x25)   # #F9A825
BLANC        = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_TEXTE   = RGBColor(0x37, 0x47, 0x4F)   # #37474F
GRIS_CLAIR   = RGBColor(0xEC, 0xEF, 0xF1)   # #ECEFF1
BLEU_INFO    = RGBColor(0x01, 0x57, 0x9B)   # #01579B
ORANGE       = RGBColor(0xE6, 0x51, 0x00)   # #E65100

W = Inches(13.33)   # largeur widescreen
H = Inches(7.5)     # hauteur widescreen


# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # entièrement vide
    return prs.slides.add_slide(layout)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=GRIS_TEXTE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def heading(slide, title, subtitle=None):
    """Bandeau titre en haut de chaque slide contenu."""
    rect(slide, 0, 0, W, Inches(1.1), fill=VERT_FONCE)
    rect(slide, 0, Inches(1.1), W, Pt(4), fill=OR)
    txbox(slide, title,
          Inches(0.35), Inches(0.1), Inches(12.5), Inches(0.75),
          size=28, bold=True, color=BLANC, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.35), Inches(0.78), Inches(12.5), Inches(0.35),
              size=14, color=VERT_CLAIR, italic=True)


def bullet_block(slide, items, x, y, w, h,
                 size=15, color=GRIS_TEXTE, spacing=Pt(6)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color


def card(slide, title, lines, x, y, w, h,
         bg=VERT_PALE, title_bg=VERT_MOYEN, title_color=BLANC,
         text_size=13, title_size=16):
    rect(slide, x, y, w, h, fill=bg, line=VERT_CLAIR, line_w=Pt(1))
    rect(slide, x, y, w, Inches(0.42), fill=title_bg)
    txbox(slide, title, x + Inches(0.1), y + Pt(4), w - Inches(0.2), Inches(0.38),
          size=title_size, bold=True, color=title_color)
    bullet_block(slide, lines,
                 x + Inches(0.15), y + Inches(0.48),
                 w - Inches(0.3), h - Inches(0.55),
                 size=text_size, color=GRIS_TEXTE, spacing=Pt(3))


def footer(slide, num, total):
    rect(slide, 0, H - Inches(0.3), W, Inches(0.3), fill=VERT_FONCE)
    txbox(slide, "Projet Gestion Palmeraie  |  UFHB – Master 2 BDGL",
          Inches(0.3), H - Inches(0.28), Inches(10), Inches(0.28),
          size=9, color=VERT_CLAIR)
    txbox(slide, f"{num} / {total}",
          Inches(12.3), H - Inches(0.28), Inches(0.8), Inches(0.28),
          size=9, color=BLANC, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

def slide_titre(prs):
    s = blank_slide(prs)
    # Fond dégradé simulé avec deux rectangles
    rect(s, 0, 0, W, H * 0.6, fill=VERT_FONCE)
    rect(s, 0, H * 0.6, W, H * 0.4, fill=VERT_PALE)
    # Bande dorée décorative
    rect(s, 0, H * 0.6 - Pt(5), W, Pt(8), fill=OR)
    # Titre principal
    txbox(s, "Gestion de Palmeraie",
          Inches(1), Inches(1.0), Inches(11), Inches(1.4),
          size=52, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    # Sous-titre
    txbox(s, "Application web de gestion agricole — Récoltes, Travaux, Ventes & Traçabilité",
          Inches(1), Inches(2.5), Inches(11), Inches(0.8),
          size=20, color=VERT_CLAIR, align=PP_ALIGN.CENTER, italic=True)
    # Séparateur
    rect(s, Inches(4), Inches(3.4), Inches(5.3), Pt(2), fill=OR)
    # Informations académiques
    infos = [
        "UFHB – Université Félix Houphouët-Boigny",
        "UFR  Sciences des Structures de la Matière et de Technologie (SSMT)",
        "Master 2 — Bases de Données et Génie Logiciel (BDGL)",
        "Mémoire de fin d'études  •  2024 – 2025",
    ]
    bullet_block(s, infos,
                 Inches(1.5), Inches(3.6), Inches(10), Inches(2.0),
                 size=16, color=GRIS_TEXTE, spacing=Pt(8))
    # Pastille technos
    techs = [("Django", VERT_MOYEN), ("React", BLEU_INFO), ("SQLite", OR), ("PWA", ORANGE)]
    for i, (label, col) in enumerate(techs):
        bx = Inches(2.8 + i * 2.0)
        rect(s, bx, Inches(6.0), Inches(1.6), Inches(0.55), fill=col)
        txbox(s, label, bx + Inches(0.05), Inches(6.03),
              Inches(1.5), Inches(0.5),
              size=16, bold=True, color=BLANC, align=PP_ALIGN.CENTER)


def slide_plan(prs, total):
    s = blank_slide(prs)
    heading(s, "Plan de la présentation")
    footer(s, 2, total)
    sections = [
        ("01", "Contexte & Objectifs",         "Problématique, périmètre fonctionnel"),
        ("02", "Choix des technologies",        "Stack backend, frontend, base de données"),
        ("03", "Architecture générale",         "Modules, interactions, déploiement"),
        ("04", "Base de données",               "Schéma, tables, relations clés"),
        ("05", "Fonctionnalités détaillées",    "Toutes les pages et flux applicatifs"),
        ("06", "Sécurité & Traçabilité",        "Audit, permissions, hors-ligne"),
        ("07", "Conclusion & Perspectives",     "Bilan et évolutions envisagées"),
    ]
    for i, (num, titre, desc) in enumerate(sections):
        bx = Inches(0.4) if i % 2 == 0 else Inches(6.8)
        by = Inches(1.35 + (i // 2) * 1.42)
        if i == 6:
            bx = Inches(3.6)
        rect(s, bx, by, Inches(6.1), Inches(1.25), fill=VERT_PALE, line=VERT_CLAIR, line_w=Pt(1))
        rect(s, bx, by, Inches(0.9), Inches(1.25), fill=VERT_MOYEN)
        txbox(s, num, bx + Inches(0.05), by + Inches(0.3), Inches(0.8), Inches(0.6),
              size=24, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        txbox(s, titre, bx + Inches(1.0), by + Inches(0.08), Inches(5.0), Inches(0.45),
              size=15, bold=True, color=VERT_FONCE)
        txbox(s, desc, bx + Inches(1.0), by + Inches(0.55), Inches(5.0), Inches(0.6),
              size=12, color=GRIS_TEXTE, italic=True)


def slide_contexte(prs, total):
    s = blank_slide(prs)
    heading(s, "Contexte & Objectifs", "Problématique agricole et périmètre du projet")
    footer(s, 3, total)

    card(s, "Problématique",
         ["• La gestion des palmeraies repose encore largement sur des fiches papier",
          "• Risques de perte de données, erreurs de calcul et absence de traçabilité",
          "• Difficulté à consolider les statistiques et à piloter la production",
          "• Aucun outil adapté au contexte terrain (connectivité intermittente)"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.4),
         text_size=13)

    card(s, "Objectifs du projet",
         ["• Numériser l'ensemble du processus de récolte et de vente",
          "• Fournir un tableau de bord décisionnel en temps réel",
          "• Assurer la traçabilité complète des modifications (audit trail)",
          "• Permettre la saisie hors ligne avec synchronisation automatique",
          "• Gérer les droits d'accès par rôle (admin / superviseur)"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.4),
         text_size=13)

    card(s, "Périmètre fonctionnel",
         ["• Fiches de récolte (régimes grands / moyens / petits)",
          "• Reçus de vente avec calcul automatique des prix",
          "• Fiches de travaux agricoles (désherbage, traitement, fertilisation…)",
          "• Gestion des récolteurs, agents terrain et secteurs",
          "• Suivi des équipements et matériels",
          "• Journal d'audit et notifications in-app"],
         Inches(0.3), Inches(3.9), Inches(12.7), Inches(2.4),
         text_size=13)


def slide_techno(prs, total):
    s = blank_slide(prs)
    heading(s, "Choix des technologies", "Justification de la stack technique")
    footer(s, 4, total)

    cards = [
        ("Backend — Django REST Framework",
         VERT_MOYEN, VERT_PALE,
         ["• Python 3.11 — maturité, richesse des librairies scientifiques",
          "• Django 4.x — ORM puissant, migrations gérées automatiquement",
          "• DRF (Django REST Framework) — API REST sécurisée avec TokenAuthentication",
          "• Architecture en applications Django indépendantes (modules)",
          "• Serializers pour validation et transformation des données"]),
        ("Frontend — React + Vite",
         BLEU_INFO, RGBColor(0xE3, 0xF2, 0xFD),
         ["• React 18 — composants réutilisables, hooks, Context API",
          "• Vite — serveur de développement rapide, build optimisé",
          "• Axios — client HTTP avec intercepteurs pour le token d'auth",
          "• CSS pur — design maîtrisé sans dépendance à un framework UI",
          "• IndexedDB (via utilitaire offline.js) pour la persistance hors ligne"]),
        ("Base de données — SQLite / PostgreSQL",
         OR, RGBColor(0xFF, 0xF8, 0xE1),
         ["• SQLite en développement — zéro configuration",
          "• Compatible PostgreSQL pour la production",
          "• 22 tables, relations clés via FK et ManyToMany",
          "• Migrations versionnées avec Django (historique complet)"]),
        ("Infrastructure & PWA",
         ORANGE, RGBColor(0xFF, 0xF3, 0xE0),
         ["• Service Worker — cache des ressources statiques (network-first)",
          "• IndexedDB — stockage des fiches hors ligne avec token intégré",
          "• Synchronisation automatique au retour de la connexion",
          "• Déploiement : Django sur serveur WSGI + React build statique"]),
    ]

    positions = [
        (Inches(0.3),  Inches(1.3)),
        (Inches(6.8),  Inches(1.3)),
        (Inches(0.3),  Inches(4.1)),
        (Inches(6.8),  Inches(4.1)),
    ]
    for (title, tc, bg, lines), (bx, by) in zip(cards, positions):
        card(s, title, lines, bx, by, Inches(6.2), Inches(2.65),
             bg=bg, title_bg=tc, text_size=12, title_size=14)


def slide_archi_generale(prs, total):
    s = blank_slide(prs)
    heading(s, "Architecture générale", "Vue d'ensemble des couches et interactions")
    footer(s, 5, total)

    # Schéma en couches
    layers = [
        ("Navigateur Web (Client)", BLEU_INFO,
         "React 18 + Vite  •  Context API  •  Axios  •  IndexedDB  •  Service Worker"),
        ("API REST (Serveur)", VERT_MOYEN,
         "Django 4 + DRF  •  TokenAuthentication  •  8 applications  •  Serializers + Views"),
        ("Base de données", OR,
         "SQLite / PostgreSQL  •  22 tables  •  ORM Django  •  Migrations versionnées"),
    ]
    arrows = ["HTTP / JSON  ↕  (CORS géré par django-cors-headers)", "ORM Django  ↕  (queries SQL)"]

    y = Inches(1.4)
    for i, (label, col, desc) in enumerate(layers):
        rect(s, Inches(1.5), y, Inches(10.3), Inches(1.0), fill=col)
        txbox(s, label, Inches(1.6), y + Pt(5),
              Inches(3.2), Inches(0.5),
              size=16, bold=True, color=BLANC)
        txbox(s, desc, Inches(4.9), y + Pt(8),
              Inches(6.8), Inches(0.5),
              size=13, color=BLANC, italic=True)
        if i < 2:
            ay = y + Inches(1.0)
            txbox(s, arrows[i],
                  Inches(4.5), ay + Pt(2), Inches(5.0), Inches(0.38),
                  size=11, color=GRIS_TEXTE, italic=True, align=PP_ALIGN.CENTER)
            rect(s, Inches(6.3), ay, Pt(2), Inches(0.42), fill=GRIS_TEXTE)
            y += Inches(1.42)
        else:
            y += Inches(1.0)

    # Modules backend
    txbox(s, "Modules applicatifs Django",
          Inches(0.3), Inches(5.1), Inches(12.7), Inches(0.35),
          size=14, bold=True, color=VERT_FONCE)
    modules = ["accounts", "recoltes", "travaux", "secteurs", "agents", "recolteurs", "materiels", "dashboard"]
    colors_m = [VERT_MOYEN, VERT_MOYEN, VERT_MOYEN, BLEU_INFO, BLEU_INFO, BLEU_INFO, OR, ORANGE]
    for i, (mod, col) in enumerate(zip(modules, colors_m)):
        bx = Inches(0.3 + i * 1.6)
        rect(s, bx, Inches(5.5), Inches(1.5), Inches(0.65), fill=col)
        txbox(s, mod, bx + Pt(4), Inches(5.55),
              Inches(1.42), Inches(0.55),
              size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)


def slide_archi_modules(prs, total):
    s = blank_slide(prs)
    heading(s, "Architecture — Modules Django", "Responsabilités de chaque application")
    footer(s, 6, total)

    modules = [
        ("accounts",    VERT_MOYEN,
         ["Authentification (login/logout)", "Gestion des utilisateurs & rôles",
          "Permissions fonctionnelles (Droit)", "Notifications in-app",
          "Réinitialisation de mot de passe", "Journal d'audit (AuditLog)"]),
        ("recoltes",    VERT_FONCE,
         ["Fiches de récolte (CRUD complet)", "Lignes et détails de récolte",
          "Reçus de vente (FicheRecuVente)", "Gestion des clients",
          "Paramètres bonus & prix officiel", "Journal d'actions (ActionLog)"]),
        ("travaux",     BLEU_INFO,
         ["Fiches de travaux agricoles", "Consommables utilisés",
          "Répartition des tâches / ouvriers", "Calcul automatique du coût total",
          "Statut d'avancement", "Lien avec secteurs et matériels"]),
        ("secteurs",    OR,
         ["Référentiel des secteurs", "Code auto-généré",
          "Superficie, relief, type de sol", "Rendement cible"]),
        ("agents",      ORANGE,
         ["Superviseurs généraux", "Agents terrain",
          "Lien avec compte utilisateur Django", "Matricule & téléphone"]),
        ("recolteurs",  RGBColor(0x6A, 0x1B, 0x9A),
         ["Personnel récolteur (alias Recolteur)", "Numéro de téléphone unique",
          "Wave / WhatsApp", "Historique des lignes de récolte"]),
        ("materiels",   RGBColor(0x00, 0x69, 0x5C),
         ["Équipements & matériels", "Table intermédiaire avec travaux",
          "Quantité utilisée par fiche", "Dates de maintenance"]),
        ("dashboard",   GRIS_TEXTE,
         ["Pas de modèle propre", "Agrégations statistiques via ORM",
          "Données consolidées pour le tableau de bord"]),
    ]

    cols = 4
    for i, (name, col, lines) in enumerate(modules):
        bx = Inches(0.25 + (i % cols) * 3.27)
        by = Inches(1.3 + (i // cols) * 2.85)
        card(s, name, lines, bx, by, Inches(3.1), Inches(2.65),
             bg=VERT_PALE, title_bg=col, text_size=11, title_size=13)


def slide_bdd_schema(prs, total):
    s = blank_slide(prs)
    heading(s, "Base de données — Vue d'ensemble", "22 tables, organisées en 5 domaines")
    footer(s, 7, total)

    domaines = [
        ("Comptes & Sécurité", VERT_MOYEN,
         ["auth_user (Django built-in)", "profil_utilisateur", "droit",
          "notification", "password_reset_token", "audit_log", "action_log"]),
        ("Récoltes & Ventes", VERT_FONCE,
         ["fiche_recolte", "ligne_recolte", "detail_recolte",
          "superviseur_adjoint", "recu_vente", "client", "parametre_bonus"]),
        ("Travaux Agricoles", BLEU_INFO,
         ["fiche_travaux", "consommable_travaux",
          "repartition_tache", "fiche_travaux_secteur"]),
        ("Référentiels Terrain", OR,
         ["secteur", "superviseur_general",
          "agent_terrain", "personnel"]),
        ("Matériels", ORANGE,
         ["materiel", "materiel_travaux"]),
    ]

    for i, (title, col, tables) in enumerate(domaines):
        bx = Inches(0.25 + (i % 3) * 4.37)
        by = Inches(1.3 + (i // 3) * 3.0)
        if i >= 3:
            bx = Inches(0.9 + (i - 3) * 5.5)
        card(s, title, ["  • " + t for t in tables],
             bx, by, Inches(4.1), Inches(2.7),
             bg=VERT_PALE, title_bg=col, text_size=12, title_size=14)


def slide_bdd_comptes(prs, total):
    s = blank_slide(prs)
    heading(s, "Base de données — Comptes & Sécurité")
    footer(s, 8, total)

    tables = [
        ("profil_utilisateur",
         ["user (OneToOne → auth_user)", "role : admin | superviseur | superviseur_adjoint",
          "must_change_password (bool)", "numero_telephone",
          "droits (M2M → droit)"]),
        ("droit",
         ["code (unique)", "label", "description",
          "ordre (tri d'affichage)",
          "→ Permissions fonctionnelles attribuées aux superviseurs"]),
        ("notification",
         ["user (FK → auth_user)", "message", "type : success | warning | info",
          "lu (bool)", "lien (URL relative)", "created_at"]),
        ("password_reset_token",
         ["user (FK)", "token (64 chars, unique)", "created_at",
          "used (bool)", "→ Valide 1 heure, usage unique"]),
        ("audit_log",
         ["utilisateur (FK → auth_user)", "table_concernee",
          "id_enregistrement", "champ_modifie",
          "ancienne_valeur / nouvelle_valeur", "motif"]),
        ("action_log",
         ["acteur (FK → auth_user)", "superviseur (FK)",
          "action (37 types possibles)", "fiche (FK optionnelle)",
          "recu (FK optionnelle)", "detail (JSON)", "annule (bool)"]),
    ]
    positions = [
        (Inches(0.25), Inches(1.3)),
        (Inches(4.55), Inches(1.3)),
        (Inches(8.85), Inches(1.3)),
        (Inches(0.25), Inches(4.05)),
        (Inches(4.55), Inches(4.05)),
        (Inches(8.85), Inches(4.05)),
    ]
    for (title, lines), (bx, by) in zip(tables, positions):
        card(s, title, lines, bx, by, Inches(4.1), Inches(2.65),
             text_size=11, title_size=13)


def slide_bdd_recoltes(prs, total):
    s = blank_slide(prs)
    heading(s, "Base de données — Récoltes & Ventes")
    footer(s, 9, total)

    tables = [
        ("fiche_recolte",
         ["statut : brouillon | soumis | valide",
          "date, superviseur_general_obj (FK)",
          "bareme_grands / moyens / petits (FCFA)",
          "depense_nourriture / transport / salaire",
          "created_by / validated_by (FK auth_user)",
          "validated_at, depense_total"]),
        ("ligne_recolte",
         ["fiche (FK → fiche_recolte)",
          "recolteur (FK → personnel)",
          "recolteur_nom (cache texte)",
          "regime_type : grands | moyens | petits",
          "salaire_calcule (auto)",
          "→ Plusieurs lignes par fiche"]),
        ("detail_recolte",
         ["ligne (FK → ligne_recolte)",
          "secteur (FK → secteur)",
          "secteur_code (cache)",
          "quantite (nbre de régimes)",
          "UNIQUE (ligne, secteur)"]),
        ("recu_vente",
         ["fiche (FK → fiche_recolte)",
          "statut : brouillon | valide",
          "client_obj (FK → client)",
          "pesee_kg, non_conformes_pct",
          "montant, prix_officiel",
          "mode_paiement, vehicule_transport"]),
        ("client",
         ["nom (unique)", "telephone", "adresse",
          "→ Acheteurs des régimes de palme"]),
        ("parametre_bonus",
         ["Singleton (pk=1)",
          "bareme_grands/moyens/petits_defaut",
          "seuil_non_conformes (%)",
          "montant_bonus (FCFA)",
          "prix_kg_officiel"]),
    ]
    positions = [
        (Inches(0.25), Inches(1.3)),
        (Inches(4.55), Inches(1.3)),
        (Inches(8.85), Inches(1.3)),
        (Inches(0.25), Inches(4.1)),
        (Inches(4.55), Inches(4.1)),
        (Inches(8.85), Inches(4.1)),
    ]
    for (title, lines), (bx, by) in zip(tables, positions):
        card(s, title, lines, bx, by, Inches(4.1), Inches(2.55),
             text_size=11, title_size=13)


def slide_bdd_travaux(prs, total):
    s = blank_slide(prs)
    heading(s, "Base de données — Travaux & Matériels")
    footer(s, 10, total)

    tables = [
        ("fiche_travaux",
         ["statut : brouillon | soumis | valide",
          "superviseur_travaux, nature_travaux",
          "superficie_couverte_ha, periode_travaux",
          "nb_personnes, observations",
          "date_debut / date_fin",
          "cout_total_calcule, salaire_total",
          "statut_avancement : planifié | en cours | terminé",
          "secteurs_couverts (M2M → secteur)"]),
        ("consommable_travaux",
         ["fiche (FK → fiche_travaux)",
          "designation, quantite, unite",
          "prix_unitaire",
          "fournisseur, numero_lot",
          "date_peremption",
          "→ sous-total = quantite × prix_unitaire"]),
        ("repartition_tache",
         ["fiche (FK → fiche_travaux)",
          "nom_prenom, nature_taches",
          "quantite, prix_unitaire",
          "salaire_total_calcule (auto via save())",
          "matricule_ouvrier"]),
        ("materiel",
         ["numero (unique)", "designation",
          "quantite, etat_physique",
          "categorie : véhicule | outil | équipement lourd | petit matériel",
          "date_acquisition, valeur_achat",
          "fournisseur, localisation, responsable",
          "date_derniere/prochaine_maintenance"]),
        ("materiel_travaux",
         ["materiel (FK)", "fiche_travaux (FK)",
          "quantite_utilisee",
          "observations",
          "UNIQUE (materiel, fiche_travaux)",
          "→ Table intermédiaire explicite"]),
    ]
    positions = [
        (Inches(0.25), Inches(1.3)),
        (Inches(5.0),  Inches(1.3)),
        (Inches(9.5),  Inches(1.3)),
        (Inches(0.25), Inches(4.15)),
        (Inches(6.9),  Inches(4.15)),
    ]
    sizes = [
        (Inches(4.6), Inches(2.9)),
        (Inches(4.3), Inches(2.9)),
        (Inches(3.6), Inches(2.9)),
        (Inches(6.4), Inches(2.65)),
        (Inches(5.9), Inches(2.65)),
    ]
    for (title, lines), (bx, by), (w, h) in zip(tables, positions, sizes):
        card(s, title, lines, bx, by, w, h, text_size=11, title_size=13)


def slide_bdd_referentiels(prs, total):
    s = blank_slide(prs)
    heading(s, "Base de données — Référentiels Terrain")
    footer(s, 11, total)

    tables = [
        ("secteur",
         ["code (unique, auto-généré depuis nom)",
          "nom, superficie_ha",
          "situation_relief, type_sol",
          "age_moyen_plants, nb_palmiers",
          "rendement_cible_t_ha",
          "statut : actif | inactif"]),
        ("personnel  (alias Recolteur)",
         ["nom, lieu_residence",
          "numero_telephone (unique)",
          "whatsapp_actif, est_wave (bools)",
          "date_naissance",
          "→ Récolteurs journaliers de la palmeraie"]),
        ("superviseur_general",
         ["user (OneToOne → auth_user, optionnel)",
          "code (auto : SUP-001, SUP-002…)",
          "nom, prenom, matricule",
          "telephone, actif",
          "→ Créé automatiquement à la création d'un compte superviseur"]),
        ("agent_terrain",
         ["nom, prenom, matricule",
          "telephone",
          "secteur (FK → secteur)",
          "actif",
          "→ Agents affectés à un secteur spécifique"]),
    ]
    positions = [
        (Inches(0.3),  Inches(1.3)),
        (Inches(3.55), Inches(1.3)),
        (Inches(6.8),  Inches(1.3)),
        (Inches(10.05),Inches(1.3)),
    ]
    for (title, lines), (bx, by) in zip(tables, positions):
        card(s, title, lines, bx, by, Inches(3.1), Inches(4.2),
             text_size=12, title_size=13)


def slide_fonc_auth(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Authentification & Rôles",
            "Sécurité basée sur les tokens DRF et les rôles métier")
    footer(s, 12, total)

    card(s, "Flux de connexion",
         ["1. L'utilisateur saisit ses identifiants sur /login",
          "2. Le backend valide via ModelBackend Django",
          "3. Un token DRF est généré et retourné",
          "4. Le token est stocké en sessionStorage (expire à la fermeture)",
          "5. Toutes les requêtes API portent l'en-tête  Authorization: Token <token>",
          "6. Si le compte est désactivé → message explicite côté frontend"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.6), text_size=13)

    card(s, "3 rôles métier",
         ["Administrateur",
          "   → Accès complet : validation, rejet, gestion des utilisateurs",
          "   → Peut modifier des fiches validées (tracé dans AuditLog)",
          "Superviseur",
          "   → Crée et soumet les fiches, gère récolteurs & secteurs",
          "   → Droits affinés par l'admin via la table droit",
          "Superviseur Adjoint",
          "   → Droits réduits, sous supervision"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.6), text_size=12)

    card(s, "Réinitialisation de mot de passe",
         ["• Clic sur « Mot de passe oublié » → saisie de l'e-mail",
          "• Backend génère un token 64 chars (secrets.token_urlsafe) valide 1 h",
          "• E-mail envoyé avec lien /reset-password?token=XXX",
          "• Le token est à usage unique et supprimé après utilisation",
          "• Les anciens tokens non utilisés sont purgés avant création"],
         Inches(0.3), Inches(4.1), Inches(6.2), Inches(2.65), text_size=13)

    card(s, "Changement de mot de passe obligatoire",
         ["• Flag must_change_password sur le profil utilisateur",
          "• Activé par défaut lors de la création d'un compte par l'admin",
          "• À la connexion, si actif → redirection forcée vers /change-password",
          "• Désactivé automatiquement après que l'utilisateur a changé son MDP",
          "• Empêche l'accès à toute autre page tant que non rempli"],
         Inches(6.7), Inches(4.1), Inches(6.3), Inches(2.65), text_size=13)


def slide_fonc_dashboard(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Tableau de bord",
            "4 onglets statistiques avec données en temps réel")
    footer(s, 13, total)

    onglets = [
        ("Vue d'ensemble", VERT_MOYEN,
         ["Nombre de fiches (brouillon / soumis / validé)",
          "Total régimes récoltés (grands + moyens + petits)",
          "Total recettes de vente",
          "Total dépenses récolte et travaux",
          "Solde net (recettes − dépenses)",
          "Graphique évolution mensuelle des récoltes"]),
        ("Secteurs", BLEU_INFO,
         ["Production par secteur (nb régimes)",
          "Carte / tableau des secteurs les plus productifs",
          "Comparaison entre secteurs",
          "Statut actif/inactif des secteurs"]),
        ("Financier", OR,
         ["Recettes vs dépenses par mois",
          "Évolution du prix officiel du kg",
          "Taux de conformité des régimes",
          "Rapport prix calculé / prix officiel"]),
        ("Récolteurs", ORANGE,
         ["Top récolteurs par nombre de régimes",
          "Salaires calculés par récolteur",
          "Historique d'activité par récolteur",
          "Statistiques de performance individuelle"]),
    ]

    for i, (titre, col, lines) in enumerate(onglets):
        bx = Inches(0.25 + i * 3.27)
        card(s, titre, lines, bx, Inches(1.3), Inches(3.1), Inches(5.8),
             bg=VERT_PALE, title_bg=col, text_size=12, title_size=14)

    txbox(s, "Note : le dashboard est accessible par les administrateurs et superviseurs. "
             "Les données sont calculées à la volée via l'ORM Django (pas de table dédiée).",
          Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35),
          size=11, color=GRIS_TEXTE, italic=True)


def slide_fonc_recoltes(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Fiches de Récolte",
            "Cycle de vie complet : brouillon → soumis → validé")
    footer(s, 14, total)

    card(s, "Création d'une fiche",
         ["• Date, superviseur général (lié à l'auth user)",
          "• Barème FCFA par type de régime (grands/moyens/petits)",
          "• Ajout de lignes : un récolteur → un type de régime",
          "• Détail par secteur : quantité de régimes récoltés",
          "• Dépenses : nourriture, transport, salaires (calculé auto)",
          "• Observations, superviseurs adjoints"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.7), text_size=13)

    card(s, "Cycle de validation",
         ["1. Brouillon   → saisi par le superviseur",
          "2. Soumis      → envoyé pour validation (lecture seule)",
          "3. Validé      → approuvé par l'admin (salaires calculés)",
          "   ou Rejeté   → retour en brouillon avec motif",
          "• Toute modification post-validation est tracée dans audit_log",
          "• Annulation possible via interface admin (action_log.annule)"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.7), text_size=13)

    card(s, "Calcul des salaires",
         ["• salaire = somme des quantités × barème du régime",
          "• Barème par défaut depuis parametre_bonus",
          "• Personnalisable par fiche (override au moment de la saisie)",
          "• depense_salaire = cumul des salaires de toutes les lignes",
          "• depense_total = nourriture + transport + salaire"],
         Inches(0.3), Inches(4.2), Inches(6.2), Inches(2.5), text_size=13)

    card(s, "Export & Consultation",
         ["• Impression de la fiche en PDF (fenêtre navigateur)",
          "• Filtrage par statut, date, superviseur",
          "• Recherche textuelle instantanée",
          "• FicheDialog : vue détaillée lecture seule",
          "• Export Excel disponible pour l'historique"],
         Inches(6.7), Inches(4.2), Inches(6.3), Inches(2.5), text_size=13)


def slide_fonc_ventes(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Reçus de Vente",
            "Traçabilité commerciale liée aux fiches de récolte")
    footer(s, 15, total)

    card(s, "Création d'un reçu",
         ["• Lié à une fiche de récolte validée",
          "• Client sélectionné dans la liste ou saisi librement",
          "• Pesée (kg), taux de non-conformes (%)",
          "• Montant saisi manuellement",
          "• Prix officiel du kg (admin uniquement)",
          "• Mode de paiement : espèce | virement",
          "• Référence facture, véhicule de transport"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.5), text_size=13)

    card(s, "Calculs automatiques",
         ["• prix_calculé = montant ÷ pesée_kg",
          "• rapport_prix = prix_calculé ÷ prix_officiel",
          "  → Indicateur qualité (>1 : favorable, <1 : défavorable)",
          "• Montant auto-calculé si prix officiel renseigné :",
          "  montant = pesée_kg × prix_officiel",
          "• Calculs effectués à la volée (propriétés Python)"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(3.5), text_size=13)

    card(s, "Validation des reçus",
         ["• Même cycle brouillon → soumis → validé que les fiches",
          "• Seul l'admin peut valider ou rejeter",
          "• Historique complet dans action_log",
          "• Affichage du rapport prix en vert (≥1) ou rouge (<1)"],
         Inches(0.3), Inches(4.95), Inches(6.2), Inches(2.2), text_size=13)

    card(s, "Paramètres Bonus (admin)",
         ["• Barèmes par défaut modifiables en un clic",
          "• Seuil de non-conformes déclenchant un bonus",
          "• Montant du bonus configurable",
          "• Prix officiel du kg (référence nationale)"],
         Inches(6.7), Inches(4.95), Inches(6.3), Inches(2.2), text_size=13)


def slide_fonc_travaux(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Fiches de Travaux",
            "Suivi des interventions agricoles avec calcul de coût")
    footer(s, 16, total)

    card(s, "Données de la fiche travaux",
         ["• Nature des travaux (10 types : désherbage, fertilisation…)",
          "• Période (AAAA-MM-JJ – AAAA-MM-JJ), validée côté serveur",
          "• Superficie couverte (ha), nombre de personnes",
          "• Secteurs couverts (sélection multiple)",
          "• Statut d'avancement : planifié | en cours | terminé",
          "• Observations"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.9), text_size=13)

    card(s, "Consommables & Main d'œuvre",
         ["Consommables :",
          "  • Désignation, quantité, unité, prix unitaire",
          "  • Fournisseur, numéro de lot, date de péremption",
          "Répartition des tâches :",
          "  • Nom de l'ouvrier, nature des tâches",
          "  • Quantité × prix unitaire = salaire",
          "Salaire total journalier saisi directement"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.9), text_size=13)

    card(s, "Calcul du coût total",
         ["coût_total = Σ(consommables) + Σ(répartitions) + salaire_total",
          "• Recalculé via recalculer_cout() côté serveur",
          "• Aussi recalculé en temps réel côté frontend",
          "• Stocké en base pour les agrégations du dashboard",
          "• UPDATE direct (sans signal) pour éviter les boucles récursives"],
         Inches(0.3), Inches(4.35), Inches(6.2), Inches(2.9), text_size=13)

    card(s, "Matériels utilisés",
         ["• Sélection depuis la liste des équipements disponibles",
          "• Quantité utilisée + observations",
          "• Table intermédiaire explicite (materiel_travaux)",
          "• Contrainte UNIQUE (matériel, fiche_travaux)",
          "• Visible dans le détail de la fiche travaux"],
         Inches(6.7), Inches(4.35), Inches(6.3), Inches(2.9), text_size=13)


def slide_fonc_recolteurs(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Gestion des Récolteurs & Agents",
            "Annuaire du personnel terrain")
    footer(s, 17, total)

    card(s, "Récolteurs (Personnel)",
         ["• Nom complet, lieu de résidence",
          "• Numéro de téléphone (identifiant unique)",
          "• Indicateurs de paiement mobile : Wave, WhatsApp",
          "• Date de naissance (validée : 1930 → aujourd'hui)",
          "• Historique complet des lignes de récolte",
          "• Fiche récolteur avec statistiques de production",
          "• Export Excel de la liste",
          "• Recherche et filtrage instantanés"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(4.5), text_size=13)

    card(s, "Agents Terrain",
         ["• Nom, prénom, matricule",
          "• Affecté à un secteur spécifique",
          "• Téléphone de contact",
          "• Statut actif / inactif",
          "• Visible dans l'annuaire agents",
          "• Lié aux superviseurs adjoints des fiches"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.4), text_size=13)

    card(s, "Superviseurs Généraux",
         ["• Créé automatiquement à la création d'un compte superviseur",
          "• Code auto-incrémenté : SUP-001, SUP-002…",
          "• Synchronisation automatique avec le compte Django",
          "• Affiché dans les fiches de récolte et travaux",
          "• Téléphone affiché dans les reçus de vente",
          "• Profil détaillé avec ses fiches et statistiques"],
         Inches(6.7), Inches(3.85), Inches(6.3), Inches(2.2), text_size=13)


def slide_fonc_secteurs(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Secteurs & Matériels",
            "Référentiels de la palmeraie")
    footer(s, 18, total)

    card(s, "Gestion des Secteurs",
         ["• Code unique auto-généré depuis le nom (slugify)",
          "  ex : « Bloc Nord » → BLOC_NORD",
          "• Superficie (ha), situation de relief, type de sol",
          "• Âge moyen des plants, nombre de palmiers",
          "• Rendement cible (t/ha)",
          "• Statut actif / inactif",
          "• Page détail : production historique, fiches liées,",
          "  agents affectés, fiches de travaux"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(4.5), text_size=13)

    card(s, "Gestion des Matériels",
         ["• Numéro unique, désignation",
          "• Catégorie : véhicule | outil | équipement lourd | petit matériel",
          "• Quantité disponible, état physique",
          "• Valeur d'achat, fournisseur",
          "• Dates de maintenance (dernière / prochaine)",
          "• Localisation, responsable",
          "• Lié aux fiches de travaux (matériels utilisés)",
          "• Vue tableau avec filtrage et tri"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(4.5), text_size=13)


def slide_fonc_audit(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Traçabilité & Audit",
            "Deux journaux complémentaires pour un audit trail complet")
    footer(s, 19, total)

    card(s, "ActionLog — Journal des actions",
         ["• 37 types d'actions couverts (création, validation, suppression…)",
          "• Acteur (qui a fait quoi), superviseur concerné",
          "• Référence à la fiche ou au reçu impacté",
          "• Détail JSON (champs avant/après, motif)",
          "• Flag annule : permet la réversion d'une action par l'admin",
          "• Visible dans Journal d'audit (admin) et Mon Audit (superviseur)",
          "• Filtrable par type, acteur, date, statut annulé"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.5), text_size=13)

    card(s, "AuditLog — Journal des modifications de champs",
         ["• Déclenché automatiquement via signal post_save Django",
          "• Enregistre chaque champ modifié sur une fiche validée",
          "• Ancienne valeur → Nouvelle valeur",
          "• Motif de modification saisi par l'admin",
          "• Différent de ActionLog : granularité au niveau du champ",
          "• Visible dans la page Audit Détail (admin uniquement)"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(3.5), text_size=13)

    card(s, "Connexions & Sécurité",
         ["• Tentatives de connexion réussies → loguées (connexion_reussie)",
          "• Tentatives échouées → loguées (tentative_connexion_echouee)",
          "• Connexion avec compte désactivé → action spécifique",
          "• Ces entrées sont visibles dans le journal admin"],
         Inches(0.3), Inches(5.0), Inches(6.2), Inches(2.2), text_size=13)

    card(s, "Annulation d'actions (admin)",
         ["• L'admin peut annuler une action superviseur",
          "• Boîte de dialogue de confirmation avec motif",
          "• L'entrée action_log est marquée  annule=True",
          "• Une nouvelle entrée annulation_action est créée",
          "• Piste d'audit complète et non destructive"],
         Inches(6.7), Inches(5.0), Inches(6.3), Inches(2.2), text_size=13)


def slide_fonc_offline(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Mode Hors Ligne (PWA)",
            "Saisie continue même sans connexion réseau")
    footer(s, 20, total)

    card(s, "Stockage hors ligne (IndexedDB)",
         ["• 2 stores : pending_recoltes  et  pending_travaux",
          "• La fiche est sérialisée avec le token DRF de l'utilisateur",
          "• Token stocké dans IndexedDB (pas sessionStorage)",
          "  → La sync réussit même si le navigateur a été fermé",
          "• savedAt timestamp pour chaque fiche",
          "• API : saveRecolteOffline(payload, token)",
          "        saveTravauxOffline(payload, token)"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.5), text_size=13)

    card(s, "Synchronisation automatique",
         ["• Déclenchée dès le retour de la connexion (event online)",
          "• syncPending() : boucle sur les items en attente",
          "  → POST direct via fetch() avec le token stocké",
          "  → Suppression de l'item si succès",
          "  → Conservation de l'item si erreur",
          "• Retourne { synced, errors, failedItems }",
          "• OfflineBanner affiche le résultat en temps réel"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(3.5), text_size=13)

    card(s, "Service Worker (mise en cache)",
         ["• Cache des ressources statiques à l'installation",
          "• Stratégie network-first pour les navigations",
          "• Stratégie cache-first pour les assets statiques",
          "• API calls toujours network-only (pas de cache API)",
          "• Enregistrement limité à la production (import.meta.env.PROD)"],
         Inches(0.3), Inches(5.0), Inches(6.2), Inches(2.2), text_size=13)

    card(s, "Bandeau OfflineBanner",
         ["• Rouge : mode hors ligne actif + nb fiches en attente",
          "• Orange : synchronisation en cours",
          "• Rouge foncé + liste d'erreurs : sync partielle",
          "• Disparaît automatiquement quand tout est sync",
          "• Sondage toutes les 5 secondes du compteur IndexedDB"],
         Inches(6.7), Inches(5.0), Inches(6.3), Inches(2.2), text_size=13)


def slide_fonc_notifs(prs, total):
    s = blank_slide(prs)
    heading(s, "Fonctionnalités — Notifications & Export",
            "Alertes in-app et extraction de données")
    footer(s, 21, total)

    card(s, "Notifications in-app (cloche)",
         ["• Cloche dans la barre de navigation (NotificationBell)",
          "• Badge rouge avec le nombre de notifications non lues",
          "• 3 types : success (vert) | warning (orange) | info (bleu)",
          "• Créées programmatiquement par le backend",
          "  → Ex : « Votre fiche du 12/06 a été validée »",
          "• Clic sur une notification → redirection vers la page concernée",
          "• Marquage comme lue individuel ou en masse",
          "• Sondage automatique toutes les 60 secondes"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(4.5), text_size=13)

    card(s, "Export Excel",
         ["Historique des récoltes :",
          "  • Fiche date, superviseur, régimes, salaires, dépenses",
          "Historique des ventes :",
          "  • Client, pesée, montant, prix officiel, rapport",
          "Liste des récolteurs :",
          "  • Nom, téléphone, résidence, Wave",
          "Liste des secteurs :",
          "  • Code, nom, superficie, statut",
          "• Format .xlsx généré côté backend (openpyxl)",
          "• Déclenché via bouton « Exporter »"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(4.5), text_size=13)


def slide_securite(prs, total):
    s = blank_slide(prs)
    heading(s, "Sécurité", "Mesures mises en place à chaque couche")
    footer(s, 22, total)

    card(s, "Authentification",
         ["• Token DRF unique par session (révocable côté serveur)",
          "• sessionStorage côté client (non accessible aux scripts inter-domaines)",
          "• Changement de MDP obligatoire à la première connexion",
          "• Tokens de reset MDP : 64 chars, 1 heure, usage unique",
          "• Verrouillage des comptes par l'admin (is_active=False)"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.6), text_size=13)

    card(s, "Contrôle d'accès",
         ["• Vérification du rôle sur chaque vue DRF (IsAuthenticated + custom)",
          "• Admin principal non modifiable par les autres admins",
          "• Droits fonctionnels granulaires pour les superviseurs",
          "• Isolation des données : superviseur ne voit que ses fiches",
          "• Frontend : routes protégées par AuthContext + PrivateRoute"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(2.6), text_size=13)

    card(s, "Validation des données",
         ["• Backend : serializers DRF avec validation custom",
          "  (dates, montants, formats, unicité téléphone…)",
          "• Frontend : validation avant soumission (dates futures,",
          "  numéros de téléphone 8-15 chiffres, valeurs négatives)",
          "• Champs numériques : sanitizeDecimal / sanitizeInt",
          "  (suppression du signe − dès la frappe)"],
         Inches(0.3), Inches(4.1), Inches(6.2), Inches(2.6), text_size=13)

    card(s, "Traçabilité & Non-répudiation",
         ["• Toute action est enregistrée dans action_log",
          "• Toute modification post-validation dans audit_log",
          "• Tentatives de connexion (réussies et échouées) loguées",
          "• Annulations conservées (non destructives) dans la table",
          "• Historique complet consultable par l'admin"],
         Inches(6.7), Inches(4.1), Inches(6.3), Inches(2.6), text_size=13)


def slide_conclusion(prs, total):
    s = blank_slide(prs)
    heading(s, "Conclusion & Perspectives")
    footer(s, total, total)

    card(s, "Bilan du projet",
         ["✔  Application web complète couvrant tout le cycle récolte → vente",
          "✔  Architecture modulaire Django + React, maintenable et évolutive",
          "✔  22 tables de base de données avec relations bien définies",
          "✔  Traçabilité complète via deux journaux d'audit complémentaires",
          "✔  Mode hors ligne fonctionnel avec synchronisation automatique",
          "✔  Sécurité multi-couches (token, rôles, validation, audit)",
          "✔  Interface responsive utilisable sur tablette terrain"],
         Inches(0.3), Inches(1.3), Inches(6.2), Inches(4.5), text_size=13)

    card(s, "Perspectives d'amélioration",
         ["→  Cartographie GPS des secteurs et des palmiers",
          "→  Application mobile native (React Native / Flutter)",
          "→  Intégration d'un module de prévision de récolte (ML)",
          "→  Rapports PDF automatiques envoyés par e-mail",
          "→  Authentification multi-facteurs (MFA)",
          "→  API publique documentée (Swagger / OpenAPI)",
          "→  Déploiement cloud (AWS / Azure) avec CI/CD"],
         Inches(6.7), Inches(1.3), Inches(6.3), Inches(4.5), text_size=13)

    rect(s, 0, H - Inches(1.15), W, Inches(0.85), fill=VERT_FONCE)
    txbox(s, "Merci pour votre attention",
          Inches(1), H - Inches(1.12), Inches(11.3), Inches(0.75),
          size=28, bold=True, color=BLANC, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# ASSEMBLAGE
# ════════════════════════════════════════════════════════════════════════════

TOTAL = 23   # nombre total de slides

def build():
    prs = new_prs()

    slide_titre(prs)                  # 1
    slide_plan(prs, TOTAL)            # 2
    slide_contexte(prs, TOTAL)        # 3
    slide_techno(prs, TOTAL)          # 4
    slide_archi_generale(prs, TOTAL)  # 5
    slide_archi_modules(prs, TOTAL)   # 6
    slide_bdd_schema(prs, TOTAL)      # 7
    slide_bdd_comptes(prs, TOTAL)     # 8
    slide_bdd_recoltes(prs, TOTAL)    # 9
    slide_bdd_travaux(prs, TOTAL)     # 10
    slide_bdd_referentiels(prs, TOTAL)# 11
    slide_fonc_auth(prs, TOTAL)       # 12
    slide_fonc_dashboard(prs, TOTAL)  # 13
    slide_fonc_recoltes(prs, TOTAL)   # 14
    slide_fonc_ventes(prs, TOTAL)     # 15
    slide_fonc_travaux(prs, TOTAL)    # 16
    slide_fonc_recolteurs(prs, TOTAL) # 17
    slide_fonc_secteurs(prs, TOTAL)   # 18
    slide_fonc_audit(prs, TOTAL)      # 19
    slide_fonc_offline(prs, TOTAL)    # 20
    slide_fonc_notifs(prs, TOTAL)     # 21
    slide_securite(prs, TOTAL)        # 22
    slide_conclusion(prs, TOTAL)      # 23

    out = "Presentation_Gestion_Palmeraie.pptx"
    prs.save(out)
    print(f"[OK]  Presentation generee : {out}  ({TOTAL} slides)")

if __name__ == "__main__":
    build()
